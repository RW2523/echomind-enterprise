#!/usr/bin/env python3
"""EchoMind pitch deck as a fully editable 16:9 PowerPoint."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import math

OUT = "/home/echomind/Documents/echomind/echomind-enterprise/marketing_doc/EchoMind_Pitch_Deck.pptx"

INK   = RGBColor(0x0B, 0x1F, 0x2A)
ACC   = RGBColor(0x0E, 0x8F, 0x9E)
TEAL  = RGBColor(0x0E, 0x8F, 0x9E)
BLUE  = RGBColor(0x1D, 0x6F, 0xB8)
PURP  = RGBColor(0x4F, 0x52, 0xC4)
AMBER = RGBColor(0xA2, 0x62, 0x0B)
GREEN = RGBColor(0x12, 0x80, 0x5F)
ROSE  = RGBColor(0xB0, 0x3A, 0x5B)
SLATE = RGBColor(0x3B, 0x4C, 0x7A)
BODY  = RGBColor(0x2B, 0x36, 0x3D)
MUTE  = RGBColor(0x6C, 0x7C, 0x86)
SOFT  = RGBColor(0x93, 0xAE, 0xB8)
WHT   = RGBColor(0xFF, 0xFF, 0xFF)
FILL  = RGBColor(0xF4, 0xF8, 0xF9)
LINE  = RGBColor(0xD3, 0xDF, 0xE4)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
SW, SH = 13.333, 7.5
M = 0.62                      # left/right margin
CW = SW - 2 * M
BLANK = prs.slide_layouts[6]


def tint(col, a):
    return RGBColor(*[int(255 - (255 - ch) * a) for ch in (col[0], col[1], col[2])])


def rect(sl, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def text(sl, x, y, w, h, runs, size=12, color=BODY, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.0):
    """runs: str or list of (text, dict-overrides)."""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        txt, ov = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.space_after = Pt(ov.get("space_after", space_after))
        p.line_spacing = ov.get("line_spacing", line_spacing)
        p.text = txt
        if ov.get("bullet"):
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(int(0.20 * 914400)))
            pPr.set("indent", str(int(-0.20 * 914400)))
            bu = etree.SubElement(pPr, qn("a:buChar"))
            bu.set("char", "\u2022")
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.italic = ov.get("italic", italic)
            r.font.color.rgb = ov.get("color", color)
    return tb


def slide_base(title, kicker="", accent=ACC, sub="", num=None):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, SW, 0.105, fill=accent)
    y = 0.52
    if kicker:
        text(sl, M, y, CW, 0.26, kicker.upper(), 11, accent, bold=True)
        y += 0.30
    text(sl, M, y, CW, 0.62, title, 33, INK, bold=True)
    y += 0.72
    if sub:
        text(sl, M, y, CW * 0.92, 0.5, sub, 13.5, MUTE, line_spacing=1.12)
        y += 0.30 + 0.22 * (1 + len(sub) // 118)
    text(sl, M, SH - 0.42, 3.0, 0.24, "EchoMind by Ajace AI", 9, MUTE)
    if num:
        text(sl, SW - M - 1.0, SH - 0.42, 1.0, 0.24, str(num), 9, MUTE, align=PP_ALIGN.RIGHT)
    return sl, y + 0.10


def bullets_box(sl, x, y, w, items, color, size=11.5, gap=5, lead=1.05):
    runs = [(t, {"bullet": True, "size": size, "color": BODY, "space_after": gap,
                 "line_spacing": lead}) for t in items]
    return text(sl, x, y, w, 0.4 * len(items) + 0.4, runs)


N = 0
def nxt():
    global N; N += 1; return N


# ═════════ 1 TITLE ═════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=INK)
rect(sl, 0, SH - 0.14, SW, 0.14, fill=ACC)
text(sl, M, 1.20, 6, 0.3, "AJACE AI", 13, ACC, bold=True)
text(sl, M, 1.70, 9, 1.2, "EchoMind", 66, WHT, bold=True)
text(sl, M, 2.95, 10, 0.5, "The private AI assistant that can sit in the room.", 21, SOFT)
text(sl, M, 3.55, 8.2, 1.0,
     "Answers from your own documents. Listens to your meetings. Speaks when spoken to. "
     "Runs entirely on one machine you own — with the network cable unplugged.",
     14, RGBColor(0x7F, 0xA6, 0xB2), line_spacing=1.25)
cx = M
for t in ["On-premises", "Cited answers", "Full-duplex voice", "Multi-tenant", "Peer-reviewed"]:
    w = 0.20 + 0.098 * len(t)
    s = rect(sl, cx, 5.95, w, 0.34, fill=INK, line=ACC, lw=1.0)
    tf = s.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.text = t; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.name = FONT; r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = ACC
    cx += w + 0.14
text(sl, M, 6.62, 4, 0.24, "Internal / Confidential", 9.5, RGBColor(0x5A, 0x72, 0x80))
nxt()

# ═════════ 2 PROBLEM ═════════
sl, y = slide_base("AI stopped at the office door", "The problem", ROSE,
                   "The organizations with the most to gain from AI are the ones least able to use it.", nxt())
items = [
    ("Confidential work is off-limits", "Legal, clinical, financial and government teams cannot paste real material into a cloud chatbot. So the AI everyone talks about is used for drafting emails, not for the work that matters."),
    ("Knowledge is scattered and silent", "Policies sit in one system, contracts in another — and the most valuable context of all, what was actually said in a meeting, is never written down."),
    ("Generic AI invents answers", "A confident, unsourced answer is not a smaller version of a correct one. In a regulated workflow it is a liability."),
    ("Expertise cannot be everywhere", "One compliance officer, one senior partner, one experienced clinician cannot sit in every room at once."),
]
colw = (CW - 0.5) / 2
for i, (t, txt_) in enumerate(items):
    px = M + (i % 2) * (colw + 0.5)
    py = y + (i // 2) * 1.42
    rect(sl, px, py, 0.36, 0.05, fill=ROSE)
    text(sl, px, py + 0.14, colw, 0.34, t, 16, INK, bold=True)
    text(sl, px, py + 0.52, colw, 0.8, txt_, 11.5, BODY, line_spacing=1.15)
by = y + 2 * 1.42 + 0.12
rect(sl, M, by, CW, 0.52, fill=tint(ROSE, .10))
text(sl, M + 0.22, by + 0.13, CW - 0.4, 0.3,
     "The result: the more sensitive and valuable your knowledge, the less AI you are allowed to use on it.",
     14, INK, bold=True)

# ═════════ 3 SOLUTION ═════════
sl, y = slide_base("Bring the AI to the data — not the data to the AI", "The solution", ACC,
                   "EchoMind is a complete AI assistant on a single appliance inside the customer's own walls.", nxt())
cards = [
    ("Reads your documents", "Upload policies, contracts, manuals. Ask in plain language. Every answer cites its source — document, section and page.", TEAL),
    ("Hears your meetings", "Live transcription with a background check against your own policy, flagging what is unsupported or non-compliant as it is said.", BLUE),
    ("Talks with you", "A real spoken conversation — responds in about half a second and stops the moment you interrupt.", PURP),
    ("Writes your documents", "Turns a brief into a finished PDF or PowerPoint in your template — minutes, memos, reports, decks.", AMBER),
    ("Keeps tenants apart", "One box can serve legal, clinical and finance teams with knowledge that cannot cross between them.", ROSE),
    ("Never phones home", "No cloud, no third-party processor, no training on your data. It works with the cable unplugged.", GREEN),
]
cw3 = (CW - 2 * 0.34) / 3
for i, (t, txt_, col) in enumerate(cards):
    px = M + (i % 3) * (cw3 + 0.34)
    py = y + (i // 3) * 1.86
    rect(sl, px, py, cw3, 1.66, fill=FILL, line=tint(col, .5), lw=1.0)
    rect(sl, px, py, cw3, 0.06, fill=col)
    text(sl, px + 0.20, py + 0.24, cw3 - 0.4, 0.3, t, 15, INK, bold=True)
    text(sl, px + 0.20, py + 0.62, cw3 - 0.4, 0.9, txt_, 11, BODY, line_spacing=1.14)

# ═════════ 4 THE IDEA ═════════
sl, y = slide_base("An assistant that can stay in the room", "The idea", PURP,
                   "Because it runs on-premises and speaks naturally, EchoMind is not confined to a browser tab. "
                   "It can be installed where the work actually happens.", nxt())
places = [
    ("IN THE MEETING ROOM", "A permanent participant", ["hears the whole discussion", "flags a wrong statement live", "writes the minutes and actions"], BLUE),
    ("ON THE DESK", "A personal assistant", ["answers from the case or chart", "drafts the document", "never sends data outside"], PURP),
    ("AT THE FRONT LINE", "A support co-pilot", ["listens to the customer call", "surfaces the right answer", "checks what was promised"], AMBER),
    ("IN THE FIELD", "An offline expert", ["works with no connectivity", "hands-free voice", "full manual library on board"], GREEN),
]
pw = (CW - 3 * 0.28) / 4
for i, (kick, t, bl, col) in enumerate(places):
    px = M + i * (pw + 0.28)
    rect(sl, px, y, pw, 2.62, fill=WHT, line=col, lw=1.6)
    rect(sl, px, y, pw, 0.48, fill=col)
    text(sl, px, y + 0.13, pw, 0.26, kick, 10.5, WHT, bold=True, align=PP_ALIGN.CENTER)
    text(sl, px, y + 0.62, pw, 0.3, t, 14.5, INK, bold=True, align=PP_ALIGN.CENTER)
    bullets_box(sl, px + 0.22, y + 1.06, pw - 0.44, bl, col, 11, 6, 1.08)
fy = y + 2.62 + 0.30
text(sl, M, fy, CW, 0.28, "What makes this possible", 12.5, ACC, bold=True)
row = [("No cloud dependency", "it can live in a locked room or a vehicle"),
       ("Speaks and listens", "it works without a keyboard, hands-free"),
       ("Knows when to speak", "silent by default; surfaces only high-confidence findings"),
       ("Separated knowledge", "the room's tenant sees only the room's content")]
cw4 = (CW - 3 * 0.26) / 4
for i, (t, s) in enumerate(row):
    px = M + i * (cw4 + 0.26)
    text(sl, px, fy + 0.34, cw4, 0.26, t, 11.5, INK, bold=True)
    text(sl, px, fy + 0.62, cw4, 0.5, s, 10.5, BODY, line_spacing=1.12)


# ═════════ 5–8 PERSONAS ═════════
def persona(kicker, title, sub, accent, who, pains, does, moment, quote):
    sl, y = slide_base(title, kicker, accent, sub, nxt())
    colw = (CW - 0.44) / 2
    rx = M + colw + 0.44
    text(sl, M, y, colw, 0.24, "WHO", 11, accent, bold=True)
    text(sl, M, y + 0.30, colw, 0.6, who, 13, BODY, line_spacing=1.2)
    py = y + 1.10
    text(sl, M, py, colw, 0.24, "WHAT SLOWS THEM DOWN TODAY", 11, accent, bold=True)
    bullets_box(sl, M, py + 0.32, colw, pains, accent, 12, 7, 1.1)
    text(sl, rx, y, colw, 0.24, "WHAT ECHOMIND DOES", 11, accent, bold=True)
    bullets_box(sl, rx, y + 0.30, colw, does, accent, 12, 7, 1.1)
    n_lines = sum(max(1, math.ceil(len(t) / 58)) for t in does)
    my = y + 0.36 + 0.255 * n_lines + 0.34
    rect(sl, rx, my, colw, 1.16, fill=tint(accent, .08), line=accent, lw=1.0)
    text(sl, rx + 0.20, my + 0.15, colw - 0.4, 0.24, "THE MOMENT IT EARNS ITS KEEP", 10, accent, bold=True)
    text(sl, rx + 0.20, my + 0.44, colw - 0.4, 0.7, moment, 11.5, BODY, line_spacing=1.14)
    qy = SH - 1.35
    rect(sl, M, qy, CW, 0.66, fill=tint(accent, .09))
    rect(sl, M, qy, 0.07, 0.66, fill=accent)
    text(sl, M + 0.28, qy + 0.18, CW - 0.6, 0.34, quote, 15, INK, italic=True)


persona("Persona 1", "The Lawyer", "Private legal associate — on the desk and in the client meeting.", PURP,
    "Partners, associates and in-house counsel working across contracts, regulation, precedent and privileged client material.",
    ["Re-reading the same agreements to find one clause",
     "Precedent buried across matters nobody can search",
     "Public AI tools forbidden — privilege and confidentiality",
     "Client calls where a commitment is made and never recorded"],
    ["Ask across the whole contract library and get the clause, with document, section and page",
     "Contract-review memos drafted to the firm's template",
     "Client meetings transcribed, with commitments and deadlines extracted",
     "Live flagging of one-sided, missing or risky clauses against the firm's playbook"],
    "Mid-negotiation, counsel asks what the liability cap was in the 2023 agreement — and gets the clause and page in seconds, without the file leaving the building.",
    "\"It reads every contract we have ever signed, and it can tell me where it got the answer.\"")

persona("Persona 2", "The Doctor", "Clinical assistant — at the desk, at the bedside, hands-free.", GREEN,
    "Physicians, nurses and clinical governance teams working under protocol, with PHI that cannot leave the institution.",
    ["Guidelines change; recalling the current one takes time",
     "Documentation eats the evening",
     "PHI rules out every cloud AI service",
     "A drug interaction can be missed under pressure"],
    ["Ask protocols and formulary in plain language, with the source shown",
     "Hands-free voice for use while examining or scrubbed in",
     "Consultations transcribed and turned into a structured visit note",
     "Live flagging of interactions, contraindications and dosing"],
    "During a consultation the assistant quietly flags that the proposed medication interacts with the patient's existing prescription — while the patient is still in the room.",
    "\"Documentation support and a second pair of eyes — and the data never leaves the hospital.\"")

persona("Persona 3", "The Manager", "Meeting facilitator — a permanent participant in the room.", BLUE,
    "Executives, chiefs of staff, PMO and operations leaders who live in meetings and own the follow-through.",
    ["Decisions are made and then disputed a month later",
     "Minutes are somebody's evening job, and arrive late",
     "Actions have no owner and no due date",
     "Nobody can search what was said last quarter"],
    ["Sits in the room and transcribes the whole discussion",
     "Separates speakers and produces an analyzed report",
     "Extracts decisions, commitments and action items with owners",
     "Flags anything said that contradicts company policy",
     "Makes every past meeting searchable alongside documents"],
    "Someone asks what was agreed about the vendor in March. The answer comes back with the sentence, the speaker and the date — and the discussion moves on.",
    "\"Minutes stop being a chore, and last quarter stops being a black hole.\"")

persona("Persona 4", "The Support Agent", "Customer care co-pilot — on the call, in real time.", AMBER,
    "Contact-center and customer-care teams answering product, policy, billing and warranty questions at speed.",
    ["New agents take months to learn the product",
     "Answers vary between agents — and some are wrong",
     "Hold time while the agent searches a wiki",
     "Promises made on calls that policy does not support"],
    ["Listens to the live call and surfaces the right answer from the real policy",
     "Every answer cited, so the agent can quote with confidence",
     "Flags a promise that contradicts warranty or pricing policy",
     "Produces the call summary and follow-up automatically",
     "Customer data never leaves the contact center"],
    "A new agent is asked an obscure warranty question. The correct answer, with the clause behind it, appears before the customer finishes the sentence.",
    "\"A first-week agent answers like a five-year veteran — and we can prove what was said.\"")

# ═════════ 9 HOW IT WORKS ═════════
sl, y = slide_base("Three steps, then it just works", "How it works", TEAL,
                   "No integration project, no data migration, no cloud account.", nxt())
steps = [("1", "Install", "One appliance, five containers, one command. No connectivity required.", TEAL),
         ("2", "Feed it", "Upload documents. Let it listen to meetings. It indexes both into one searchable knowledge base.", BLUE),
         ("3", "Use it", "Type, talk, or leave it running in the room. Every answer cites its source.", GREEN)]
sw3 = (CW - 2 * 0.44) / 3
for i, (n, t, txt_, col) in enumerate(steps):
    px = M + i * (sw3 + 0.44)
    rect(sl, px, y, 0.62, 0.62, fill=col, shape=MSO_SHAPE.OVAL)
    text(sl, px, y + 0.15, 0.62, 0.34, n, 22, WHT, bold=True, align=PP_ALIGN.CENTER)
    text(sl, px + 0.78, y + 0.10, sw3 - 0.8, 0.36, t, 19, INK, bold=True)
    text(sl, px, y + 0.82, sw3, 0.8, txt_, 12.5, BODY, line_spacing=1.16)
uy = y + 1.92
text(sl, M, uy, CW, 0.28, "Under the hood — for the technical buyer", 13, SLATE, bold=True)
under = [("Hybrid retrieval", "Meaning search and exact-keyword search run together, then a second model re-reads the best passages against the question."),
         ("Grounded generation", "The answer model sees only the retrieved evidence, fenced as data — so a document cannot give it instructions."),
         ("Honest refusal", "If nothing relevant clears the gate, it says the answer is not in the documents rather than guessing."),
         ("Isolation in the query", "The tenant check runs inside the search itself, so another tenant's content never enters the ranking.")]
cw2 = (CW - 0.44) / 2
for i, (t, s) in enumerate(under):
    px = M + (i % 2) * (cw2 + 0.44)
    py = uy + 0.40 + (i // 2) * 0.86
    rect(sl, px, py + 0.07, 0.07, 0.07, fill=SLATE)
    text(sl, px + 0.18, py, cw2 - 0.2, 0.26, t, 12, INK, bold=True)
    text(sl, px + 0.18, py + 0.26, cw2 - 0.2, 0.52, s, 10.8, BODY, line_spacing=1.12)

# ═════════ 10 PROOF ═════════
sl, y = slide_base("Measured, not claimed", "Proof", GREEN,
                   "An instrumented evaluation of this system has been accepted for publication at an "
                   "international peer-reviewed conference (QASC 2026).", nxt())
stats = [("0.98", "citation precision"), ("78%", "correct refusals on unanswerable questions"),
         ("0 / 50", "cross-tenant leaks"), ("0.6 s", "to first spoken response"),
         ("2×", "reduction in prompt-injection success")]
sw5 = CW / 5
for i, (big, lab) in enumerate(stats):
    px = M + i * sw5
    text(sl, px, y, sw5 - 0.2, 0.6, big, 36, GREEN, bold=True)
    text(sl, px, y + 0.62, sw5 - 0.24, 0.6, lab, 11, MUTE, line_spacing=1.1)
py = y + 1.46
colw = (CW - 0.44) / 2
text(sl, M, py, colw, 0.24, "WHAT WE PROVED", 11, GREEN, bold=True)
bullets_box(sl, M, py + 0.32, colw, [
    "Answers point at the right source, and it refuses when it should.",
    "Tenant isolation held on every retrieval path audited.",
    "A structural defense halved prompt-injection success versus merely instructing the model.",
    "The dual-loop design is faster AND safer — removing the safety rule produced ungrounded claims on every turn.",
], GREEN, 12, 7, 1.1)
text(sl, M + colw + 0.44, py, colw, 0.24, "WHAT WE OPENLY REPORT AS OPEN", 11, ROSE, bold=True)
bullets_box(sl, M + colw + 0.44, py + 0.32, colw, [
    "Response timing can still reveal that a restricted document exists.",
    "Blocked injection attempts are contained but not yet surfaced to the operator.",
    "Human-panel evaluation of answer quality is not yet complete.",
], ROSE, 12, 7, 1.1)
by = SH - 1.24
rect(sl, M, by, CW, 0.60, fill=tint(GREEN, .09))
text(sl, M + 0.24, by + 0.17, CW - 0.5, 0.3,
     "Publishing the limitations is the point. Technical buyers verify claims — and almost no vendor in this category has any.",
     14, INK, bold=True)

# ═════════ 11 COMPETITION ═════════
sl, y = slide_base("Why EchoMind wins", "Competition", SLATE,
                   "Every alternative forces a trade-off between capability and control. We remove it.", nxt())
hdr = ["", "Runs offline", "Cites sources", "Refuses when unsure", "Live meeting check", "Tenant isolation"]
rows = [["Public AI chatbots", "No", "No", "No", "No", "No"],
        ["Cloud enterprise AI", "No", "Partial", "Partial", "No", "Partial"],
        ["Enterprise search", "Yes", "Links only", "n/a", "No", "Partial"],
        ["Meeting transcription", "No", "No", "n/a", "Transcript only", "No"],
        ["EchoMind", "Yes", "Yes", "Yes", "Yes", "Yes, in-query"]]
wds = [2.9, 1.62, 1.62, 1.82, 1.92, 1.82]
hh = 0.60
cx = M
for h, wd in zip(hdr, wds):
    rect(sl, cx, y, wd, hh, fill=SLATE)
    text(sl, cx + 0.14, y + 0.11, wd - 0.24, 0.40, h, 10.5, WHT, bold=True, line_spacing=1.05)
    cx += wd
ry = y + hh
for ri, row in enumerate(rows):
    last = ri == len(rows) - 1
    rh = 0.50
    rect(sl, M, ry, sum(wds), rh, fill=tint(ACC, .13) if last else (FILL if ri % 2 == 0 else WHT), line=LINE, lw=0.5)
    cx = M
    for ci, (v, wd) in enumerate(zip(row, wds)):
        if ci == 0:
            col, bold = INK, True
        else:
            good = v.startswith("Yes")
            col = GREEN if good else (AMBER if v.startswith("Partial") else MUTE)
            bold = good
        text(sl, cx + 0.14, ry + 0.14, wd - 0.2, 0.26, v, 11.5 if ci == 0 else 11, col, bold=bold)
        cx += wd
    ry += rh
ry += 0.28
text(sl, M, ry, CW, 0.28, "And one more thing nobody else has", 13, ACC, bold=True)
text(sl, M, ry + 0.32, CW * 0.8, 0.5,
     "A peer-reviewed measurement of the architecture, with the limitations published. "
     "In a market of unverifiable claims, evidence is the differentiator.", 12, BODY, line_spacing=1.14)

# ═════════ 12 COMMERCIAL ═════════
sl, y = slide_base("How it is sold and deployed", "Commercial", SLATE,
                   "One platform, five packaged products, four deployment shapes.", nxt())
colw = (CW - 0.5) / 2
text(sl, M, y, colw, 0.24, "PACKAGED PRODUCTS", 11, SLATE, bold=True)
prods = [("EchoMind Health", "hospitals, clinics, clinical governance"),
         ("EchoMind Law", "firms and in-house legal teams"),
         ("EchoMind Bank", "advice, compliance, KYC/AML"),
         ("EchoMind Meeting Rooms", "boards, PMO, operations"),
         ("EchoMind Retail", "stores, sales floors, customer care")]
for i, (t, s) in enumerate(prods):
    py = y + 0.36 + i * 0.36
    text(sl, M, py, 2.5, 0.26, t, 12, INK, bold=True)
    text(sl, M + 2.6, py, colw - 2.6, 0.26, s, 11, MUTE)
rx = M + colw + 0.5
text(sl, rx, y, colw, 0.24, "DEPLOYMENT SHAPES", 11, SLATE, bold=True)
deps = [("On-premises", "the standard sale — an appliance in their rack"),
        ("Edge / field", "remote sites, clinics, vehicles, vessels"),
        ("Private cloud", "their own controlled tenancy"),
        ("Air-gapped", "defense, government, critical infrastructure")]
for i, (t, s) in enumerate(deps):
    py = y + 0.36 + i * 0.36
    text(sl, rx, py, 1.7, 0.26, t, 12, INK, bold=True)
    text(sl, rx + 1.8, py, colw - 1.8, 0.26, s, 11, MUTE)
cy = y + 2.30
text(sl, M, cy, CW, 0.24, "COMMERCIAL SHAPE", 11, SLATE, bold=True)
com = [("Owned hardware", "One appliance per deployment. No per-token cost that grows with usage."),
       ("No processor agreement", "Nothing is sub-processed, so the procurement path is dramatically shorter."),
       ("Expansion", "More departments = more tenants on the same box; more sites = more boxes."),
       ("Services", "Vertical content packs, template customization, and integration work.")]
cw4 = (CW - 3 * 0.28) / 4
for i, (t, s) in enumerate(com):
    px = M + i * (cw4 + 0.28)
    rect(sl, px, cy + 0.32, cw4, 1.20, fill=FILL, line=tint(SLATE, .4), lw=1.0)
    text(sl, px + 0.18, cy + 0.48, cw4 - 0.36, 0.26, t, 12, INK, bold=True)
    text(sl, px + 0.18, cy + 0.78, cw4 - 0.36, 0.66, s, 10.5, BODY, line_spacing=1.12)

# ═════════ 13 ROADMAP / ASK ═════════
sl, y = slide_base("Where it goes next", "Roadmap and ask", ACC, "", nxt())
colw = (CW - 0.5) / 2
text(sl, M, y, colw, 0.24, "SHIPPING TODAY", 11, GREEN, bold=True)
bullets_box(sl, M, y + 0.32, colw, [
    "Knowledge Chat, Live Transcript, Silent Assistant, Voice, Document Studio, Boardroom",
    "Five vertical products with enforced tenant isolation",
    "Secure Export Gateway; role-based access and audit",
], GREEN, 12, 8, 1.12)
text(sl, M + colw + 0.5, y, colw, 0.24, "NEXT", 11, AMBER, bold=True)
bullets_box(sl, M + colw + 0.5, y + 0.32, colw, [
    "Surface blocked injection attempts to the operator",
    "Constant-time responses to close the timing channel",
    "Human-panel evaluation and larger-corpus scale study",
    "Deeper meeting-room hardware integration",
], AMBER, 12, 8, 1.12)
ay = y + 2.30
rect(sl, M, ay, CW, 1.92, fill=INK)
rect(sl, M, ay, 0.09, 1.92, fill=ACC)
text(sl, M + 0.44, ay + 0.26, 4, 0.26, "The ask", 12, ACC, bold=True)
text(sl, M + 0.44, ay + 0.62, CW - 1, 0.44, "Give us one room, one department, and one week.", 24, WHT, bold=True)
text(sl, M + 0.44, ay + 1.18, CW - 1, 0.6,
     "Bring your own documents. We install, index, and demonstrate on your content — "
     "then unplug the network cable and run the whole demo again.", 13, SOFT, line_spacing=1.2)

prs.save(OUT)
print("wrote", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
