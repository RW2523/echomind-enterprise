"""PPTX extraction: text + structure (slides)."""
from __future__ import annotations
from io import BytesIO
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches


def extract_pptx(data: bytes) -> Tuple[str, List[dict]]:
    """
    Extract full text and per-slide structure.
    Returns (full_text, [{"slide": 1, "text": "..."}, ...]).
    """
    prs = Presentation(BytesIO(data))
    slides = []
    full_parts = []
    for i, s in enumerate(prs.slides):
        parts = []
        for sh in s.shapes:
            if hasattr(sh, "text") and sh.text:
                parts.append(sh.text.strip())
        text = "\n".join(parts)
        slides.append({"slide": i + 1, "text": text})
        full_parts.append(text)
    full_text = "\n\n".join(full_parts)
    return full_text, slides
