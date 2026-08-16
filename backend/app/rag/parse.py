"""
PDF / DOCX / PPTX parsing.

Primary PDF parser: PyMuPDF (fitz) for layout-aware extraction with header/footer
removal. Falls back to pypdf when fitz is unavailable.

Public API unchanged: parse_any returns (filetype, text, estimated_pages, page_offsets).
"""
from __future__ import annotations

import logging
import re
from collections import Counter
from io import BytesIO
from typing import Dict, List, Optional, Tuple

from docx import Document
from pptx import Presentation

from .normalize import normalize_extracted_text

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Header / footer deduplication
# ---------------------------------------------------------------------------

_MIN_PAGES_FOR_HF_DETECT = 6
_HF_SAMPLE_PAGES = 30
_HF_MIN_FREQ_RATIO = 0.55


def _normalize_hf_line(line: str) -> str:
    """Normalize a candidate header/footer line for frequency comparison.

    Strips page numbers, dates, and whitespace so that lines that differ only
    in the page number are counted as the same repeated header/footer.
    """
    t = line.strip()
    t = re.sub(r"\b\d{1,5}\b", "#", t)
    # Collapse masked-number chains ("#.#", "#-#", "# / #") to one "#": 'DoD 7000.14-R'
    # otherwise normalizes to 'DoD #.#-R' on some pages and 'DoD #-R' on others, splitting
    # the frequency count so neither variant crossed the detection threshold.
    t = re.sub(r"#(?:\s*[.\-–—/]\s*#)+", "#", t)
    t = re.sub(r"\s+", " ", t).strip()
    return t


def detect_header_footer_patterns(
    pages_text: List[str],
    n_lines: int = 4,
) -> Tuple[List[str], List[str]]:
    """Detect repeated header and footer lines across pages.

    Examines the first and last `n_lines` of each page. Lines that appear
    (after normalization) on >55% of sampled pages are considered headers/footers.

    Returns (header_patterns, footer_patterns) — normalized strings to match.
    """
    if len(pages_text) < _MIN_PAGES_FOR_HF_DETECT:
        return ([], [])

    sample = pages_text[:_HF_SAMPLE_PAGES]
    n_sample = len(sample)
    header_counter: Counter = Counter()
    footer_counter: Counter = Counter()

    for page_text in sample:
        lines = [l.strip() for l in page_text.split("\n") if l.strip()]
        if not lines:
            continue
        top = lines[:n_lines]
        bottom = lines[-n_lines:] if len(lines) > n_lines else []
        for line in top:
            norm = _normalize_hf_line(line)
            # len >= 3 (was 5): DoD-style page numbers like "2-45" mask to "#-#" (3 chars)
            # and were silently exempt from detection, so they survived into every chunk.
            if len(norm) >= 3:
                header_counter[norm] += 1
        for line in bottom:
            norm = _normalize_hf_line(line)
            if len(norm) >= 3:
                footer_counter[norm] += 1

    threshold = int(n_sample * _HF_MIN_FREQ_RATIO)
    headers = [pat for pat, cnt in header_counter.items() if cnt >= threshold]
    footers = [pat for pat, cnt in footer_counter.items() if cnt >= threshold]
    if headers or footers:
        logger.info(
            "parse: detected %d header pattern(s), %d footer pattern(s) from %d pages",
            len(headers), len(footers), n_sample,
        )
    return (headers, footers)


# A line that is nothing but a page marker: bare number ("45"), roman numeral,
# DoD-style compound page number ("2-45", "02a-3"), or a change-marker like
# "* June 2017". These are stripped at the outermost content lines of every page
# even when frequency detection missed them (short docs, <55%-frequency variants).
_PAGE_MARKER_RE = re.compile(
    r"^\s*(?:"
    r"\d{1,4}"                                  # 45
    r"|[ivxlcdm]{1,7}"                          # xiv
    r"|\d{1,3}[A-Za-z]?-\d{1,4}[a-z]?"          # 2-45, 02a-3
    r"|\*?\s*(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}"
    r")\s*\*?\s*$",
    re.IGNORECASE,
)


def strip_header_footer_lines(
    page_text: str,
    header_patterns: List[str],
    footer_patterns: List[str],
    n_lines: int = 4,
) -> str:
    """Remove header/footer lines from a single page's text.

    Windows are measured in CONTENT (non-empty) lines. The old version indexed raw
    lines, but block extraction joins blocks with blank lines, so a 3-line header
    occupied raw indices 0..4 and its tail escaped the window — measured on the FMR
    corpus: 86/159 pages kept their 'DoD 7000.14-R' header, and every page kept its
    '2-45'-style page number (also exempted by the old 5-char detection filter)."""
    lines = page_text.split("\n")
    if not lines:
        return page_text

    patterns = set(header_patterns + footer_patterns)
    content_idx = [i for i, l in enumerate(lines) if l.strip()]
    top_set = set(content_idx[:n_lines])
    bottom_set = set(content_idx[-n_lines:]) if content_idx else set()
    edge_set = set(content_idx[:2]) | set(content_idx[-2:])  # outermost 2 content lines

    clean: List[str] = []
    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            clean.append(line)
            continue
        if (i in top_set or i in bottom_set) and _normalize_hf_line(stripped) in patterns:
            continue
        if i in edge_set and _PAGE_MARKER_RE.match(stripped):
            continue
        clean.append(line)
    return "\n".join(clean)


# ---------------------------------------------------------------------------
# PyMuPDF (fitz) parser — primary
# ---------------------------------------------------------------------------

def _parse_pdf_fitz(data: bytes) -> Tuple[str, int, List[Tuple[int, int]]]:
    """Parse PDF using PyMuPDF for layout-aware text extraction.

    Returns (text, page_count, page_char_offsets).
    """
    import fitz  # PyMuPDF

    doc = fitz.open(stream=data, filetype="pdf")
    pages_text: List[str] = []
    for page in doc:
        # Extract by BLOCKS (≈ paragraphs), joined with a blank line, so paragraph
        # structure survives normalization. get_text("text") emits only single "\n"
        # line breaks; normalize_whitespace_preserve_paragraphs then collapses a whole
        # page into ONE line, leaving zero "\n\n" in the document. That silently made
        # _looks_like_long_form bail at `paragraph_breaks < 5`, so every PDF was typed
        # USER/FAQ and the entire BookRAG path (parent/child chunks, section index,
        # TOC routing, contextual headers) never ran. Measured on the FMR corpus:
        # 0 breaks before -> 74-206 after. sort=True keeps reading order.
        blocks = page.get_text("blocks", sort=True) or []
        parts = [
            (b[4] or "").strip()
            for b in blocks
            if len(b) > 6 and b[6] == 0 and (b[4] or "").strip()  # b[6]==0 -> text block
        ]
        pages_text.append("\n\n".join(parts) if parts else (page.get_text("text") or ""))

    header_pats, footer_pats = detect_header_footer_patterns(pages_text)

    # Normalize each page BEFORE computing offsets so page_offsets live in the same
    # coordinate system the chunker uses (which no longer re-normalizes). Otherwise
    # downstream normalization shifts every offset and page citations drift. (H3)
    page_offsets: List[Tuple[int, int]] = []
    cleaned_pages: List[str] = []
    offset = 0
    for i, raw_text in enumerate(pages_text):
        cleaned = normalize_extracted_text(strip_header_footer_lines(raw_text, header_pats, footer_pats))
        page_offsets.append((offset, i + 1))
        cleaned_pages.append(cleaned)
        # +2 must match the 2-char "\n\n" page separator below, or every page_offset
        # drifts by 1 char per page and page citations point at the wrong page (H3).
        offset += len(cleaned) + 2

    text = "\n\n".join(cleaned_pages)
    doc.close()
    return text, len(pages_text), page_offsets


# ---------------------------------------------------------------------------
# pypdf parser — fallback
# ---------------------------------------------------------------------------

def _parse_pdf_pypdf(data: bytes) -> Tuple[str, int, List[Tuple[int, int]]]:
    """Fallback PDF parser using pypdf (text-layer only, no layout awareness)."""
    from pypdf import PdfReader

    r = PdfReader(BytesIO(data))
    pages_text: List[str] = []
    for page in r.pages:
        pages_text.append(page.extract_text() or "")

    header_pats, footer_pats = detect_header_footer_patterns(pages_text)

    # Normalize per-page before computing offsets (see H3 note in _parse_pdf_fitz).
    page_offsets: List[Tuple[int, int]] = []
    cleaned_pages: List[str] = []
    offset = 0
    for i, raw_text in enumerate(pages_text):
        cleaned = normalize_extracted_text(strip_header_footer_lines(raw_text, header_pats, footer_pats))
        page_offsets.append((offset, i + 1))
        cleaned_pages.append(cleaned)
        # +2 must match the 2-char "\n\n" page separator below, or every page_offset
        # drifts by 1 char per page and page citations point at the wrong page (H3).
        offset += len(cleaned) + 2

    text = "\n\n".join(cleaned_pages)
    return text, len(r.pages), page_offsets


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def parse_pdf(data: bytes) -> Tuple[str, int, List[Tuple[int, int]]]:
    """Return (text, page_count, page_char_offsets).

    Uses PyMuPDF (fitz) for layout-aware extraction with header/footer removal.
    Falls back to pypdf if fitz is not installed.

    page_char_offsets: list of (start_char_offset, page_number_1indexed) sorted by
    ascending offset. Allows O(log n) lookup of page number for any character offset.
    """
    try:
        return _parse_pdf_fitz(data)
    except ImportError:
        logger.info("parse: PyMuPDF not available, falling back to pypdf")
        return _parse_pdf_pypdf(data)
    except Exception as exc:
        logger.warning("parse: PyMuPDF failed (%s), falling back to pypdf", exc)
        return _parse_pdf_pypdf(data)


def page_for_offset(offset: int, page_offsets: List[Tuple[int, int]]) -> Optional[int]:
    """Binary-search page_offsets to find the 1-indexed page number for a character offset.

    page_offsets must be sorted ascending by start_char_offset (guaranteed by parse_pdf).
    Returns None when page_offsets is empty.
    """
    if not page_offsets:
        return None
    lo, hi = 0, len(page_offsets) - 1
    result = page_offsets[0][1]
    while lo <= hi:
        mid = (lo + hi) // 2
        if page_offsets[mid][0] <= offset:
            result = page_offsets[mid][1]
            lo = mid + 1
        else:
            hi = mid - 1
    return result


def parse_docx(data: bytes) -> str:
    doc = Document(BytesIO(data))
    return "\n".join([p.text for p in doc.paragraphs])


def parse_pptx(data: bytes) -> str:
    prs = Presentation(BytesIO(data))
    parts = []
    for s in prs.slides:
        for sh in s.shapes:
            if hasattr(sh, "text") and sh.text:
                parts.append(sh.text)
    return "\n".join(parts)


def parse_any(filename: str, data: bytes) -> Tuple[str, str, int, List[Tuple[int, int]]]:
    """Return (filetype, text, estimated_pages, page_char_offsets).

    estimated_pages=0 and page_char_offsets=[] when page info is unavailable.
    """
    f = filename.lower()
    if f.endswith(".pdf"):
        text, pages, page_offsets = parse_pdf(data)
        return "pdf", text, pages, page_offsets
    if f.endswith(".docx"):
        return "docx", parse_docx(data), 0, []
    if f.endswith(".pptx"):
        return "pptx", parse_pptx(data), 0, []
    return "txt", data.decode("utf-8", errors="ignore"), 0, []
