from __future__ import annotations

import io
from typing import Any, Dict, List

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt


def _wrap_lines(text: str, width: int = 92) -> List[str]:
    words = (text or "").split()
    if not words:
        return [""]
    lines: List[str] = []
    cur: List[str] = []
    for w in words:
        if sum(len(x) for x in cur) + len(cur) + len(w) > width and cur:
            lines.append(" ".join(cur))
            cur = [w]
        else:
            cur.append(w)
    if cur:
        lines.append(" ".join(cur))
    return lines


def export_report_pdf(report: Dict[str, Any]) -> bytes:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.multi_cell(0, 10, report.get("title") or "Board Room Report")
    pdf.ln(4)
    pdf.set_font("Helvetica", "", 11)
    meta = []
    if report.get("session_name"):
        meta.append(f"Session: {report['session_name']}")
    if report.get("session_location"):
        meta.append(f"Location: {report['session_location']}")
    if meta:
        pdf.multi_cell(0, 6, " · ".join(meta))
        pdf.ln(2)

    summary = (report.get("executive_summary") or "").strip()
    if summary:
        pdf.set_font("Helvetica", "B", 13)
        pdf.cell(0, 8, "Executive summary", ln=True)
        pdf.set_font("Helvetica", "", 11)
        for line in _wrap_lines(summary):
            pdf.multi_cell(0, 6, line)
        pdf.ln(2)

    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 8, "Polished minutes", ln=True)
    pdf.set_font("Helvetica", "", 10)
    polished = (report.get("polished_transcript") or "").strip()
    for line in _wrap_lines(polished, 96):
        pdf.multi_cell(0, 5, line)

    checks = report.get("knowledge_checks") or []
    if checks:
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 13)
        pdf.cell(0, 8, "Knowledge validation", ln=True)
        pdf.set_font("Helvetica", "", 10)
        for i, c in enumerate(checks, 1):
            cls = str(c.get("classification") or "related").replace("_", " ").title()
            pdf.set_font("Helvetica", "B", 11)
            pdf.multi_cell(0, 6, f"{i}. {cls}")
            pdf.set_font("Helvetica", "", 10)
            pdf.multi_cell(0, 5, f"Claim: {c.get('claim', '')}")
            if c.get("interpretation"):
                pdf.multi_cell(0, 5, f"Interpretation: {c.get('interpretation')}")
            if c.get("suggested_action"):
                pdf.multi_cell(0, 5, f"Action: {c.get('suggested_action')}")
            ev = c.get("evidence") or []
            for e in ev[:2]:
                pdf.multi_cell(0, 5, f"Source: {e.get('source_name')} — {str(e.get('matched_text') or '')[:180]}")
            pdf.ln(2)

    out = pdf.output(dest="S")
    if isinstance(out, str):
        return out.encode("latin-1", errors="replace")
    return bytes(out)


def export_report_pptx(report: Dict[str, Any]) -> bytes:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = report.get("title") or "Board Room Report"
    subtitle = title_slide.placeholders[1]
    bits = []
    if report.get("session_name"):
        bits.append(report["session_name"])
    if report.get("session_location"):
        bits.append(report["session_location"])
    subtitle.text = " · ".join(bits) if bits else "EchoMind Board Room"

    summary = (report.get("executive_summary") or "").strip()
    if summary:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive summary"
        body = slide.placeholders[1].text_frame
        body.text = summary

    polished = (report.get("polished_transcript") or "").strip()
    if polished:
        chunks = []
        cur = ""
        for para in polished.split("\n"):
            if len(cur) + len(para) > 2200 and cur:
                chunks.append(cur)
                cur = para
            else:
                cur = f"{cur}\n{para}" if cur else para
        if cur:
            chunks.append(cur)
        for idx, chunk in enumerate(chunks[:8], 1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Polished minutes" if idx == 1 else f"Minutes (cont. {idx})"
            slide.placeholders[1].text_frame.text = chunk[:3000]

    checks = report.get("knowledge_checks") or []
    if checks:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Knowledge validation"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, c in enumerate(checks[:12], 1):
            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            p.text = f"{i}. [{c.get('classification')}] {c.get('claim', '')[:220]}"
            p.level = 0
            p.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
