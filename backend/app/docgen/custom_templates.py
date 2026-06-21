"""
Custom (user-uploaded) templates for Document Studio.

Two responsibilities:

1. ``parse_template_file(filename, data)`` — turn an uploaded file (PPTX / DOCX /
   PDF / MD / TXT) into a Document Studio *template blueprint* shaped exactly like
   the built-in entries in :mod:`app.docgen.templates`. We derive the section list
   from the file's structure (slide titles, Word "Heading*" paragraphs, PDF/MD/TXT
   heading-like lines) so the generation pipeline can plan content against the
   uploaded document's outline. The returned dict carries ``custom: True`` and a
   ``source_format`` tag so callers can distinguish it from a built-in template.

2. ``render_into_pptx_template(template_bytes, doc)`` — open the uploaded ``.pptx``
   with python-pptx so its theme / master / colours are PRESERVED, then APPEND
   freshly generated content slides (title + one slide per heading with
   paragraph / bullets / table / flow / callout / image content) using the
   uploaded deck's own slide layouts. It is best-effort and theme-preserving: on
   ANY failure it RAISES so the caller can fall back to the normal
   :func:`app.docgen.render_pptx.render_pptx`.

Robustness contract:
  * ``parse_template_file`` NEVER raises. Missing optional libs, corrupt files, or
    an unrecognised format all degrade to three generic sections (Overview /
    Details / Summary) so the feature always produces a usable template.
  * ``render_into_pptx_template`` is the opposite: it RAISES on any problem so the
    caller's normal renderer can take over.

No new pip dependencies — only python-pptx, Pillow, and the existing parse helpers.
"""
from __future__ import annotations

import logging
import os
import re
import base64
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

from . import theme as _theme

logger = logging.getLogger(__name__)

# ── caps so a giant uploaded file can't produce a runaway template / deck ────
_MAX_SECTIONS = 20
_MIN_SECTIONS = 3
_MAX_TITLE_CHARS = 120
_MAX_HEADING_LINE_CHARS = 90      # a "heading-like" line stays short
_MAX_SCAN_LINES = 4000            # don't scan unbounded PDF/MD/TXT bodies

# Generic fallback used whenever we cannot detect real structure.
_GENERIC_SECTIONS: Tuple[Tuple[str, str], ...] = (
    ("overview", "Overview"),
    ("details", "Details"),
    ("summary", "Summary"),
)


# ===========================================================================
# parse_template_file
# ===========================================================================
def _ext(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower().lstrip(".")


def _source_format(filename: str) -> str:
    e = _ext(filename)
    if e in ("pptx", "docx", "pdf", "md", "txt"):
        return e
    if e in ("markdown", "mdown", "mkd"):
        return "md"
    if e == "ppt":
        return "pptx"
    if e == "doc":
        return "docx"
    return "txt"


def _slugify(name: str) -> str:
    """Filesystem-name -> stable template id slug."""
    stem = os.path.splitext(os.path.basename(name or ""))[0]
    slug = re.sub(r"[^a-z0-9]+", "_", stem.lower()).strip("_")
    slug = re.sub(r"_+", "_", slug)
    if not slug:
        slug = "custom_template"
    return ("custom_" + slug)[:80]


def _doc_type_from_file(filename: str) -> str:
    """Human-friendly doc type derived from the file name stem."""
    stem = os.path.splitext(os.path.basename(filename or ""))[0]
    stem = re.sub(r"[_\-]+", " ", stem).strip()
    stem = re.sub(r"\s+", " ", stem)
    if not stem:
        return "Custom Document"
    # Title-case but keep existing all-caps acronyms.
    words = [w if (w.isupper() and len(w) > 1) else w.capitalize() for w in stem.split(" ")]
    return (" ".join(words))[:120]


def _pick_theme(filename: str, doc_type: str) -> str:
    """Pick a sensible built-in theme name from cues in the file/doc type."""
    hay = f"{filename} {doc_type}".lower()
    rules = (
        (("legal", "compliance", "contract", "policy", "regulat", "counsel"), "counsel"),
        (("pitch", "deck", "market", "sales", "case study", "case_study", "campaign"), "spotlight"),
        (("sop", "process", "procedure", "training", "guide", "learning", "onboard"), "evergreen"),
        (("business", "executive", "report", "proposal", "finance", "board", "prd", "product"), "executive"),
        (("technical", "architecture", "engineering", "system", "spec", "whitepaper", "design"), "midnight"),
    )
    available = set(_theme._THEMES.keys())  # type: ignore[attr-defined]
    for needles, name in rules:
        if any(n in hay for n in needles) and name in available:
            return name
    return _theme.DEFAULT_THEME


def _clean_title(s: Any) -> str:
    t = str(s or "").strip()
    t = re.sub(r"\s+", " ", t)
    return t[:_MAX_TITLE_CHARS]


def _section_key(title: str, used: set, idx: int) -> str:
    base = re.sub(r"[^a-z0-9]+", "_", title.lower()).strip("_")
    base = re.sub(r"_+", "_", base)[:40]
    if not base:
        base = f"section_{idx + 1}"
    key = base
    n = 2
    while key in used:
        key = f"{base}_{n}"
        n += 1
    used.add(key)
    return key


def _build_sections(titles: List[str]) -> List[Dict[str, Any]]:
    """Turn a list of detected section titles into section_blueprint entries."""
    # De-dupe consecutive/exact repeats while preserving order, drop empties.
    seen_titles: set = set()
    clean: List[str] = []
    for t in titles:
        ct = _clean_title(t)
        if not ct:
            continue
        low = ct.lower()
        if low in seen_titles:
            continue
        seen_titles.add(low)
        clean.append(ct)
        if len(clean) >= _MAX_SECTIONS:
            break

    if len(clean) < _MIN_SECTIONS:
        # Pad with generic sections (avoiding title clashes) up to the minimum.
        for _key, gtitle in _GENERIC_SECTIONS:
            if len(clean) >= _MIN_SECTIONS:
                break
            if gtitle.lower() not in seen_titles:
                clean.append(gtitle)
                seen_titles.add(gtitle.lower())

    used_keys: set = set()
    sections: List[Dict[str, Any]] = []
    for i, title in enumerate(clean):
        sections.append({
            "key": _section_key(title, used_keys, i),
            "title": title,
            "kinds": ["prose", "bullets", "table"],
            "guidance": f'Write the "{title}" section based on the brief and source.',
        })
    return sections


def _looks_like_heading(line: str) -> bool:
    """Heuristic: is this PDF/TXT line a section heading?

    Headings tend to be short, not end with sentence punctuation, and are either
    numbered (``1.``, ``2.3``, ``IV.``) or Title-Case / UPPER-CASE.
    """
    s = line.strip()
    if not s or len(s) > _MAX_HEADING_LINE_CHARS:
        return False
    if s.endswith((".", ",", ";", ":")) and not re.match(r"^\s*\d", s):
        # trailing sentence punctuation -> prose (numbered lines may end with ':')
        if s.endswith((".", ",", ";")):
            return False
    words = s.split()
    if len(words) > 14:
        return False
    # Numbered / lettered section markers.
    if re.match(r"^\s*(\d+(\.\d+)*[.)]?|[IVXLC]+[.)]|[A-Z][.)])\s+\S", s):
        return True
    # ALL CAPS short line (allow digits/punctuation).
    letters = [c for c in s if c.isalpha()]
    if letters and all(c.isupper() for c in letters) and len(words) <= 10:
        return True
    # Title Case: most significant words start uppercase, no terminal period.
    sig = [w for w in words if len(w) > 3]
    if sig and sum(1 for w in sig if w[:1].isupper()) >= max(1, int(len(sig) * 0.6)):
        return not s.endswith(".")
    return False


def _titles_from_pptx(data: bytes) -> List[str]:
    titles: List[str] = []
    try:
        from pptx import Presentation  # local import: tolerate missing lib
    except Exception as e:  # pragma: no cover - environment dependent
        logger.warning("custom_templates: python-pptx unavailable for parse: %s", e)
        return titles
    try:
        prs = Presentation(BytesIO(data))
    except Exception as e:
        logger.warning("custom_templates: could not open uploaded pptx: %s", e)
        return titles
    for slide in prs.slides:
        title = ""
        # 1) The title placeholder, if the layout has one.
        try:
            ph_title = slide.shapes.title
            if ph_title is not None and (ph_title.text or "").strip():
                title = ph_title.text
        except Exception:
            title = ""
        # 2) Otherwise the first non-empty text run on the slide.
        if not title:
            for sh in slide.shapes:
                try:
                    txt = (sh.text or "").strip() if hasattr(sh, "text") else ""
                except Exception:
                    txt = ""
                if txt:
                    title = txt.split("\n", 1)[0]
                    break
        ct = _clean_title(title)
        if ct:
            titles.append(ct)
        if len(titles) >= _MAX_SECTIONS:
            break
    return titles


def _titles_from_docx(data: bytes) -> List[str]:
    titles: List[str] = []
    try:
        from docx import Document  # local import: tolerate missing lib
    except Exception as e:  # pragma: no cover
        logger.warning("custom_templates: python-docx unavailable for parse: %s", e)
        return titles
    try:
        document = Document(BytesIO(data))
    except Exception as e:
        logger.warning("custom_templates: could not open uploaded docx: %s", e)
        return titles
    for para in document.paragraphs:
        try:
            style_name = (para.style.name or "") if para.style else ""
        except Exception:
            style_name = ""
        if style_name.strip().lower().startswith("heading"):
            ct = _clean_title(para.text)
            if ct:
                titles.append(ct)
        if len(titles) >= _MAX_SECTIONS:
            break
    return titles


def _titles_from_pdf(data: bytes) -> List[str]:
    titles: List[str] = []
    try:
        from ..rag.parse import parse_pdf  # local import: heavy deps
    except Exception as e:  # pragma: no cover
        logger.warning("custom_templates: rag.parse unavailable for pdf: %s", e)
        return titles
    try:
        text, _pages, _offsets = parse_pdf(data)
    except Exception as e:
        logger.warning("custom_templates: parse_pdf failed: %s", e)
        return titles
    for line in (text or "").split("\n")[:_MAX_SCAN_LINES]:
        if _looks_like_heading(line):
            ct = _clean_title(line)
            if ct:
                titles.append(ct)
        if len(titles) >= _MAX_SECTIONS:
            break
    return titles


def _titles_from_markdown(text: str) -> List[str]:
    titles: List[str] = []
    for raw in text.split("\n")[:_MAX_SCAN_LINES]:
        m = re.match(r"^\s{0,3}(#{1,3})\s+(.+?)\s*#*\s*$", raw)
        if m:
            ct = _clean_title(m.group(2))
            if ct:
                titles.append(ct)
        if len(titles) >= _MAX_SECTIONS:
            break
    return titles


def _titles_from_text(text: str) -> List[str]:
    """Plain text: short non-empty lines that look like headings."""
    titles: List[str] = []
    for raw in text.split("\n")[:_MAX_SCAN_LINES]:
        if _looks_like_heading(raw):
            ct = _clean_title(raw)
            if ct:
                titles.append(ct)
        if len(titles) >= _MAX_SECTIONS:
            break
    return titles


def parse_template_file(filename: str, data: bytes) -> Dict[str, Any]:
    """Parse an uploaded template file into a Document Studio template dict.

    Returns a dict shaped like a :mod:`app.docgen.templates` entry, with extra
    ``custom`` and ``source_format`` keys. NEVER raises — on any problem it
    degrades to three generic sections so the feature always works.
    """
    filename = filename or "template"
    data = data or b""
    src = _source_format(filename)
    doc_type = _doc_type_from_file(filename)

    titles: List[str] = []
    try:
        if src == "pptx":
            titles = _titles_from_pptx(data)
        elif src == "docx":
            titles = _titles_from_docx(data)
        elif src == "pdf":
            titles = _titles_from_pdf(data)
        elif src == "md":
            text = data.decode("utf-8", errors="ignore")
            titles = _titles_from_markdown(text)
            if not titles:
                titles = _titles_from_text(text)
        else:  # txt / unknown
            text = data.decode("utf-8", errors="ignore")
            titles = _titles_from_text(text)
    except Exception as e:
        logger.warning("custom_templates: parse_template_file failed for %s: %s", filename, e)
        titles = []

    sections = _build_sections(titles)

    template_id = _slugify(filename)
    name = doc_type if doc_type else "Custom Template"
    chosen_theme = _pick_theme(filename, doc_type)

    template: Dict[str, Any] = {
        "id": template_id,
        "name": name,
        "icon": "File",
        "persona": "General Assistant",
        "description": f"Custom template derived from the uploaded {src.upper()} file '{os.path.basename(filename)}'.",
        "default_org": "EchoMind",
        "default_doc_type": doc_type,
        "theme": chosen_theme,
        "system_guidance": (
            "Write a clear, well-structured professional document. Use plain, precise language, "
            "logical headings, concise paragraphs, bullet lists for enumerations, and tables for "
            "comparisons or structured data. Ground every claim in the provided brief and source "
            "material; do not invent facts."
        ),
        "section_blueprint": sections,
        "appendix": None,
        "custom": True,
        "source_format": src,
    }
    logger.info(
        "custom_templates: parsed %s (%s) -> %d sections, theme=%s",
        filename, src, len(sections), chosen_theme,
    )
    return template


# ===========================================================================
# render_into_pptx_template
# ===========================================================================
def _hex_rgb(hexstr: str):
    from pptx.dml.color import RGBColor
    r, g, b = _theme.hex_to_rgb(hexstr)
    return RGBColor(r, g, b)


def _decode_data_url(data_url: str) -> Optional[bytes]:
    """Decode a ``data:image/...;base64,...`` URL to raw bytes. None on failure."""
    try:
        if not isinstance(data_url, str) or not data_url.startswith("data:image"):
            return None
        b64 = data_url.split(",", 1)[1] if "," in data_url else ""
        if not b64:
            return None
        return base64.b64decode(b64)
    except Exception as e:
        logger.warning("custom_templates: bad image data_url: %s", e)
        return None


def _pick_layouts(prs) -> Tuple[Any, Any]:
    """Return (title_content_layout, blank_layout) from the uploaded deck.

    Chooses a Title+Content layout (a layout that has both a title placeholder and
    at least one body/content placeholder) if available, else the blankest layout
    we can find, else the first layout. Always returns usable layouts.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER

    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("uploaded pptx has no slide layouts")

    title_content = None
    blank = None
    fewest = None
    fewest_n = 10_000
    for lay in layouts:
        has_title = False
        has_body = False
        n_ph = 0
        for ph in lay.placeholders:
            n_ph += 1
            try:
                pht = ph.placeholder_format.type
            except Exception:
                pht = None
            if pht in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                has_title = True
            elif pht in (
                PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT,
                PP_PLACEHOLDER.SUBTITLE, getattr(PP_PLACEHOLDER, "CONTENT", PP_PLACEHOLDER.OBJECT),
            ):
                has_body = True
        if has_title and has_body and title_content is None:
            title_content = lay
        if n_ph == 0 and blank is None:
            blank = lay
        if n_ph < fewest_n:
            fewest_n = n_ph
            fewest = lay

    blank = blank or fewest or layouts[-1]
    title_content = title_content or blank or layouts[0]
    return title_content, blank


def _add_title(slide, text: str, accent, ink) -> bool:
    """Set the slide's title placeholder if present. Returns True if used a placeholder."""
    try:
        title_ph = slide.shapes.title
    except Exception:
        title_ph = None
    if title_ph is not None:
        try:
            title_ph.text = text[:_MAX_TITLE_CHARS]
            return True
        except Exception:
            pass
    return False


def _add_title_textbox(slide, text: str, accent, prs):
    """Fallback title as a textbox near the top, in the theme accent colour."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    w = prs.slide_width
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), w - Inches(1.2), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text[:_MAX_TITLE_CHARS]
    run.font.size = Pt(30)
    run.font.bold = True
    try:
        run.font.color.rgb = accent
    except Exception:
        pass
    return box


def _body_placeholder(slide):
    """Find a usable body/content placeholder on the slide (not the title). None if absent."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    for ph in slide.placeholders:
        try:
            pht = ph.placeholder_format.type
        except Exception:
            pht = None
        if pht in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            continue
        if ph.has_text_frame:
            return ph
    return None


def _content_text_frame(slide, prs):
    """Return a text frame for body content: the body placeholder or a new textbox."""
    from pptx.util import Inches
    ph = _body_placeholder(slide)
    if ph is not None and ph.has_text_frame:
        tf = ph.text_frame
        tf.clear()
        return tf
    w = prs.slide_width
    h = prs.slide_height
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), w - Inches(1.2), h - Inches(2.2))
    box.text_frame.word_wrap = True
    return box.text_frame


def _add_para(tf, text: str, ink, *, size=16, bullet=False, ordered_n: Optional[int] = None, first_holder=[True]):  # noqa: B006
    from pptx.util import Pt
    text = (text or "").strip()
    if not text:
        return
    if not tf.paragraphs[0].runs and first_holder[0]:
        p = tf.paragraphs[0]
        first_holder[0] = False
    else:
        p = tf.add_paragraph()
    p.space_after = Pt(6)
    run = p.add_run()
    if ordered_n is not None:
        run.text = f"{ordered_n}. {text}"
    elif bullet:
        run.text = f"•  {text}"
    else:
        run.text = text
    run.font.size = Pt(size)
    try:
        run.font.color.rgb = ink
    except Exception:
        pass


def render_into_pptx_template(template_bytes: bytes, doc: Dict[str, Any]) -> bytes:
    """Render ``doc`` into the uploaded ``.pptx`` template, preserving its theme.

    Opens the uploaded deck (so its master / theme / colours stay), then appends a
    title slide and one content slide per heading using the deck's OWN layouts.
    Raises on ANY failure so the caller can fall back to the standard renderer.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    if not isinstance(doc, dict):
        raise ValueError("doc must be a dict")
    if not template_bytes:
        raise ValueError("empty template_bytes")

    prs = Presentation(BytesIO(template_bytes))  # raises on a corrupt file
    title_layout, blank_layout = _pick_layouts(prs)

    pal = _theme.get_theme(doc.get("theme"))
    accent = _hex_rgb(pal.get("primary", "0EA5B7"))
    ink = _hex_rgb(pal.get("ink", "0F172A"))
    body_color = _hex_rgb(pal.get("body", "1E293B"))
    muted = _hex_rgb(pal.get("muted", "64748B"))

    def _t(v: Any, limit: int = 4000) -> str:
        s = "" if v is None else str(v)
        s = s.strip()
        return s[:limit]

    # ── Title slide ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(title_layout)
    used_ph = _add_title(slide, _t(doc.get("title") or "Untitled Document", 200), accent, ink)
    if not used_ph:
        _add_title_textbox(slide, _t(doc.get("title") or "Untitled Document", 200), accent, prs)
    subtitle = _t(doc.get("subtitle") or "", 240)
    org = _t(doc.get("org") or "", 120)
    doc_type = _t(doc.get("doc_type") or "", 120)
    meta = "   ·   ".join([p for p in (org, doc_type) if p])
    # Put subtitle/meta into the body placeholder if one exists, else a textbox.
    sub_lines = [x for x in (subtitle, meta) if x]
    if sub_lines:
        body = _body_placeholder(slide)
        if body is not None and body.has_text_frame:
            tf = body.text_frame
            tf.clear()
            holder = [True]
            for line in sub_lines:
                _add_para(tf, line, body_color, size=18, first_holder=holder)
        else:
            box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), prs.slide_width - Inches(1.2), Inches(1.6))
            box.text_frame.word_wrap = True
            holder = [True]
            for line in sub_lines:
                _add_para(box.text_frame, line, body_color, size=18, first_holder=holder)

    # ── Content slides ─────────────────────────────────────────────────────────
    blocks = doc.get("blocks") if isinstance(doc.get("blocks"), list) else []

    cur_slide = None
    cur_tf = None
    cur_holder = [True]

    def _new_content_slide(title_text: str):
        nonlocal cur_slide, cur_tf, cur_holder
        s = prs.slides.add_slide(title_layout)
        if not _add_title(s, _t(title_text, _MAX_TITLE_CHARS), accent, ink):
            _add_title_textbox(s, _t(title_text, _MAX_TITLE_CHARS), accent, prs)
        cur_slide = s
        cur_tf = None
        cur_holder = [True]
        return s

    def _ensure_tf():
        nonlocal cur_slide, cur_tf, cur_holder
        if cur_slide is None:
            _new_content_slide("Details")
        if cur_tf is None:
            cur_tf = _content_text_frame(cur_slide, prs)
            cur_holder = [True]
        return cur_tf

    for block in blocks:
        if not isinstance(block, dict):
            continue
        btype = str(block.get("type", "")).strip().lower()
        try:
            if btype == "heading":
                text = _t(block.get("text") or "", _MAX_TITLE_CHARS)
                if text:
                    _new_content_slide(text)

            elif btype == "paragraph":
                _add_para(_ensure_tf(), _t(block.get("text") or "", 1200), body_color, size=16, first_holder=cur_holder)

            elif btype == "bullets":
                items = [x for x in (block.get("items") or []) if _t(x)]
                ordered = bool(block.get("ordered"))
                tf = _ensure_tf()
                for i, it in enumerate(items):
                    _add_para(tf, _t(it, 600), body_color, size=15,
                              bullet=not ordered, ordered_n=(i + 1) if ordered else None,
                              first_holder=cur_holder)

            elif btype == "table":
                _render_table_slide(prs, _new_content_slide(_t(block.get("caption") or "Table", _MAX_TITLE_CHARS)),
                                    block, accent, ink, muted, body_color)
                cur_slide = None  # table consumes its own slide

            elif btype == "flow":
                tf = _ensure_tf()
                title = _t(block.get("title") or "", 200)
                if title:
                    _add_para(tf, title, accent, size=18, first_holder=cur_holder)
                steps = [x for x in (block.get("steps") or []) if _t(x)]
                for i, step in enumerate(steps):
                    suffix = "  ↓" if i < len(steps) - 1 else ""
                    _add_para(tf, f"{i + 1}. {_t(step, 300)}{suffix}", body_color, size=15, first_holder=cur_holder)

            elif btype == "callout":
                _render_callout(prs, _ensure_tf(), block, accent, body_color, cur_holder)

            elif btype == "image":
                _render_image(prs, _ensure_tf, _new_content_slide, block, muted)

            elif btype in ("divider", "pagebreak"):
                cur_slide = None
                cur_tf = None
            # unknown block types are skipped
        except Exception as e:
            # In best-effort render we keep going; the whole-call try lives in the caller.
            logger.warning("custom_templates: skipped %s block in template render: %s", btype, e)
            continue

    buf = BytesIO()
    prs.save(buf)  # raises if the deck got into a bad state
    return buf.getvalue()


def _render_table_slide(prs, slide, block, accent, ink, muted, body_color):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    cols = [str(c)[:80] for c in (block.get("columns") or [])][:6]
    rows = [r for r in (block.get("rows") or []) if isinstance(r, (list, tuple))]
    if not cols and rows:
        cols = [f"Column {i + 1}" for i in range(min(6, max(len(r) for r in rows)))]
    if not cols:
        return
    n_cols = len(cols)
    rows = rows[:12]
    norm: List[List[str]] = []
    for r in rows:
        cells = [str(c)[:120] for c in list(r)[:n_cols]]
        cells += [""] * (n_cols - len(cells))
        norm.append(cells)

    n_rows = len(norm) + 1
    w = prs.slide_width - Inches(1.2)
    h = min(prs.slide_height - Inches(2.2), Inches(0.5 * n_rows + 0.2))
    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.6), w, h)
    table = gfx.table
    for c, col in enumerate(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        tf = cell.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = col
        run.font.bold = True
        run.font.size = Pt(11)
        try:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        except Exception:
            pass
    for ri, row in enumerate(norm, start=1):
        for c in range(n_cols):
            cell = table.cell(ri, c)
            tf = cell.text_frame
            run = tf.paragraphs[0].add_run()
            run.text = row[c] if c < len(row) else ""
            run.font.size = Pt(10)
            try:
                run.font.color.rgb = body_color
            except Exception:
                pass


def _render_callout(prs, tf, block, accent, body_color, holder):  # noqa: ARG001 (prs kept for symmetry)
    style = str(block.get("style", "info")).strip().lower()
    title = str(block.get("title") or "").strip()
    text = str(block.get("text") or "").strip()
    acc_hex, _soft = _theme.CALLOUT_COLORS.get(style, _theme.CALLOUT_COLORS["info"])
    acc_color = _hex_rgb(acc_hex)
    if title:
        _add_para(tf, title, acc_color, size=16, first_holder=holder)
    if text:
        _add_para(tf, text, body_color, size=14, first_holder=holder)


def _render_image(prs, ensure_tf, new_slide, block, muted):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    caption = str(block.get("caption") or block.get("alt") or "").strip()
    img_bytes = _decode_data_url(block.get("data_url") or "")
    if img_bytes:
        slide = new_slide(caption[:_MAX_TITLE_CHARS] or "Image")
        try:
            stream = BytesIO(img_bytes)
            left = Inches(1.0)
            top = Inches(1.7)
            width = prs.slide_width - Inches(2.0)
            slide.shapes.add_picture(stream, left, top, width=width)
            return
        except Exception as e:
            logger.warning("custom_templates: add_picture failed, using placeholder: %s", e)
    # Placeholder: a muted note in the running body so layout stays intentional.
    tf = ensure_tf()
    label = caption or (str(block.get("prompt") or "").strip()[:120]) or "Image"
    _add_para(tf, f"[ Image: {label} ]", muted, size=14)
