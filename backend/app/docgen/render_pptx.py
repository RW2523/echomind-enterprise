"""
PPTX renderer for Document Studio documents (v2 — professional designed deck).

Turns a normalized document dict (see ``models.py``) into a polished, theme-driven
16:9 PowerPoint deck and returns the file as bytes.

    from .render_pptx import render_pptx
    pptx_bytes = render_pptx(doc)

Design language (driven entirely by app.docgen.theme.get_theme(doc["theme"])):
  * TITLE slide   — cover_bg background, optional cover hero image (full-bleed or
                    right-half), big cover_ink title, subtitle, and an "org · doc_type
                    · date" meta line in cover_accent.
  * AGENDA slide  — lists the level-1 (and nested level-2) headings.
  * SECTION slide — one per level-1 heading: a big number + the title on cover_bg.
  * CONTENT slide — a top accent BAR + the heading title, a FOOTER (org left / slide
                    number right). paragraph/bullets stream into the body and spill to
                    a fresh slide once it fills.
  * table  -> own content slide with a native styled pptx table (primary header fill,
              white text, zebra body, capped & truncated).
  * flow   -> own content slide: steps as a chain of rounded rectangles + arrows
              (vertical, or horizontal when few/short).
  * callout -> a colored rounded rectangle (CALLOUT_COLORS) on the current/own slide.
  * image  -> own content slide: decoded base64 picture scaled to fit + caption, or a
              tasteful placeholder rounded rect when no data_url is present.
  * CLOSING "Thank you" slide on cover_bg.

The renderer is defensive: every dynamic value is capped, malformed blocks are skipped,
images never crash the export (placeholder/skip), and at minimum the title + closing
slides are produced so a bad LLM response can never raise. Hard slide cap ~140.
"""
from __future__ import annotations

import base64
import logging
import math
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import theme as _theme
from . import models
from ..utils.ids import now_iso

logger = logging.getLogger(__name__)

# ── Slide geometry (inches), 16:9 ───────────────────────────────────────────
_SLIDE_W = 13.33
_SLIDE_H = 7.5
_MARGIN = 0.7
_BODY_LEFT = _MARGIN
_BODY_TOP = 1.55                       # below the title bar
_BODY_WIDTH = _SLIDE_W - 2 * _MARGIN
_FOOTER_TOP = _SLIDE_H - 0.5
_BODY_BOTTOM = _FOOTER_TOP - 0.15      # body must stop above the footer
_BODY_HEIGHT = _BODY_BOTTOM - _BODY_TOP

# ── Hard caps so a runaway doc can't produce a giant / broken deck ───────────
_MAX_SLIDES = 140
_MAX_PARA_CHARS = 600
_MAX_TITLE_CHARS = 130
_MAX_AGENDA_ITEMS = 22
_MAX_TABLE_ROWS = 12         # incl. header
_MAX_TABLE_COLS = 6
_MAX_CELL_CHARS = 200
_MAX_FLOW_STEPS = 8
_MAX_STEP_CHARS = 110
_MAX_CALLOUT_CHARS = 520
_MAX_CAPTION_CHARS = 200

# ── Vertical rhythm for the unified content-flow layout ──────────────────────
_PARA_GAP = 0.12          # gap below a paragraph (inches)
_BLOCK_GAP = 0.18         # gap below a callout / block element
_TB_INSET = 0.1           # python-pptx default textbox L/R inset (inches)


def _est_lines(text: Any, width_in: float, font_pt: float) -> int:
    """Estimate how many wrapped lines `text` needs in a `width_in` box at `font_pt`.
    Deliberately a slight over-estimate so blocks never overlap."""
    s = "" if text is None else str(text)
    if not s:
        return 1
    char_w = max(0.045, font_pt * 0.0099)        # ~0.158in per char at 16pt
    cpl = max(6, int(max(0.5, width_in) / char_w))
    lines = 0
    for seg in s.split("\n"):
        seg = seg.rstrip()
        lines += max(1, math.ceil(len(seg) / cpl)) if seg else 1
    return lines


def _est_text_height(text: Any, width_in: float, font_pt: float, leading: float = 1.22) -> float:
    """Estimated rendered height (inches) of wrapped text."""
    return _est_lines(text, width_in, font_pt) * (font_pt * leading / 72.0)


def _txt(v: Any, limit: int) -> str:
    """Coerce to a trimmed, length-capped string (ellipsised if truncated)."""
    if v is None:
        return ""
    if isinstance(v, dict):
        # Flatten {'title','text'}-style step/item dicts instead of leaking a Python repr.
        title = str(v.get("title") or v.get("step") or v.get("name") or v.get("label") or "").strip()
        body = str(v.get("text") or v.get("description") or v.get("detail") or v.get("value") or "").strip()
        v = f"{title}: {body}" if title and body else (title or body or
            " — ".join(str(x).strip() for x in v.values() if isinstance(x, (str, int, float)) and str(x).strip()))
    elif not isinstance(v, str):
        v = str(v)
    v = v.strip()
    if "**" in v or "__" in v:   # strip leaked markdown bold markers
        v = v.replace("**", "").replace("__", "")
    if len(v) > limit:
        cut = v[: max(1, limit - 1)]
        sp = cut.rfind(" ")
        cut = cut[:sp] if sp > limit * 0.6 else cut   # avoid mid-word cuts
        v = cut.rstrip(" ,;:-") + "…"
    return v


def _rgb(hexstr: str) -> RGBColor:
    r, g, b = _theme.hex_to_rgb(hexstr)
    return RGBColor(r, g, b)


def _decode_data_url(data_url: Any) -> Optional[bytes]:
    """Decode a 'data:image/...;base64,XXXX' string to bytes. None on any problem."""
    if not isinstance(data_url, str) or not data_url.startswith("data:image"):
        return None
    try:
        b64 = data_url.split(",", 1)[1] if "," in data_url else ""
        if not b64:
            return None
        return base64.b64decode(b64, validate=False)
    except Exception:
        return None


def _image_size(data: bytes) -> Optional[Tuple[int, int]]:
    """Return (w, h) px of the image, or None if Pillow can't read it."""
    try:
        from PIL import Image  # Pillow is installed
        with Image.open(BytesIO(data)) as im:
            return int(im.width), int(im.height)
    except Exception:
        return None


def _flyer_extract(doc: Dict[str, Any]) -> Tuple[str, List[str], str, str]:
    """(tagline, bullets, cta_title, cta_text) from a flyer document's blocks."""
    blocks = [b for b in (doc.get("blocks") or []) if isinstance(b, dict)]
    tagline = str(doc.get("subtitle") or "").strip()
    bullets: List[str] = []
    cta_title = cta_text = ""
    paras: List[str] = []
    for b in blocks:
        t = str(b.get("type", "")).lower()
        if t == "bullets" and not bullets:
            bullets = [str(x).strip() for x in (b.get("items") or []) if str(x).strip()][:5]
        elif t == "callout" and not (cta_title or cta_text):
            cta_title = str(b.get("title") or "").strip()
            cta_text = str(b.get("text") or "").strip()
        elif t == "paragraph":
            tx = str(b.get("text") or "").strip()
            if tx:
                paras.append(tx)
    if not tagline and paras:
        tagline = paras[0]
    if not bullets and len(paras) > 1:
        bullets = paras[1:6]
    if not (cta_title or cta_text):
        cta_title = "Get in touch"
        cta_text = str(doc.get("org") or "")
    return tagline, bullets, cta_title, cta_text


class _Deck:
    """Builds a themed Presentation, tracking the current content slide for streaming text."""

    def __init__(self, pal: Dict[str, str]) -> None:
        self.pal = pal
        # Resolved theme colors.
        self.c_primary = _rgb(pal["primary"])
        self.c_primary_dark = _rgb(pal["primary_dark"])
        self.c_accent = _rgb(pal["accent"])
        _ah = "".join(c for c in str(pal.get("accent") or "").lstrip("#") if c in "0123456789abcdefABCDEF")
        self._accent_hex = (_ah[:6] if len(_ah) >= 6 else "22D3EE").upper()
        self.c_ink = _rgb(pal["ink"])
        self.c_body = _rgb(pal["body"])
        self.c_muted = _rgb(pal["muted"])
        self.c_bg = _rgb(pal["bg"])
        self.c_panel = _rgb(pal["panel"])
        self.c_panel_alt = _rgb(pal["panel_alt"])
        self.c_cover_bg = _rgb(pal["cover_bg"])
        self.c_cover_ink = _rgb(pal["cover_ink"])
        self.c_cover_accent = _rgb(pal["cover_accent"])
        self.c_white = RGBColor(0xFF, 0xFF, 0xFF)

        self.org = ""

        self.prs = Presentation()
        self.prs.slide_width = Inches(_SLIDE_W)
        self.prs.slide_height = Inches(_SLIDE_H)
        self._blank = self._pick_blank_layout()

        # Unified content-flow state for the active content slide.
        self.current_section: str = ""  # most recent LEVEL-1 section (for dividers)
        self._last_heading: str = ""    # most recent heading of ANY level (title fallback)
        self._content_slide = None      # slide currently accepting body content (or None)
        self._content_y = _BODY_TOP     # next free y on the active content slide
        self._current_title = ""        # title shown on the active content slide
        self._pending_title = ""        # a heading awaiting its first body block
        self._section_index = 0         # running level-1 section counter (for dividers)

    def _pick_blank_layout(self):
        prs = self.prs
        if len(prs.slide_layouts) > 6:
            return prs.slide_layouts[6]   # canonical "Blank" layout
        return prs.slide_layouts[-1]

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    def full(self) -> bool:
        return self.slide_count >= _MAX_SLIDES

    # ── slide / shape primitives ────────────────────────────────────────────
    def _new_slide(self, bg: Optional[RGBColor] = None):
        slide = self.prs.slides.add_slide(self._blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg if bg is not None else self.c_bg
        for ph in list(slide.placeholders):
            ph.element.getparent().remove(ph.element)
        return slide

    def _rect(self, slide, shape_type, left, top, width, height, *,
              fill=None, line=None, line_w=1.0):
        shp = slide.shapes.add_shape(
            shape_type, Inches(left), Inches(top), Inches(width), Inches(height),
        )
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(line_w)
        shp.shadow.inherit = False
        return shp

    def _textbox(self, slide, text, left, top, width, height, *,
                 size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 anchor=None, font=None):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color if color is not None else self.c_body
        if font:
            run.font.name = font
        return box

    def _footer(self, slide):
        """Org (left) + slide number (right) at the bottom of a content slide."""
        if self.org:
            self._textbox(
                slide, _txt(self.org, 80),
                _MARGIN, _FOOTER_TOP, _SLIDE_W / 2 - _MARGIN, 0.35,
                size=9, color=self.c_muted, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE,
            )
        self._textbox(
            slide, str(self.slide_count),
            _SLIDE_W / 2, _FOOTER_TOP, _SLIDE_W / 2 - _MARGIN, 0.35,
            size=9, color=self.c_muted, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    def _title_bar(self, slide, title):
        """Top accent bar + heading title for a content slide."""
        title = _txt(title, _MAX_TITLE_CHARS)
        # Thin full-width accent bar.
        self._rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, _SLIDE_W, 0.16, fill=self.c_primary)
        # A short thicker accent tab under the title for emphasis.
        self._rect(slide, MSO_SHAPE.RECTANGLE, _MARGIN, 1.28, 2.4, 0.05, fill=self.c_accent)
        # Shrink long titles so they never wrap into / past the accent tab below.
        tlen = len(title)
        size = 28 if tlen <= 46 else (23 if tlen <= 72 else 19)
        self._textbox(
            slide, title,
            _MARGIN, 0.42, _SLIDE_W - 2 * _MARGIN, 0.8,
            size=size, bold=True, color=self.c_primary_dark,
            anchor=MSO_ANCHOR.MIDDLE, font=self.pal.get("font_bold"),
        )

    def _reset_body(self):
        self._content_slide = None
        self._content_y = _BODY_TOP
        self._current_title = ""
        self._pending_title = ""

    # ── unified content-flow helpers ──────────────────────────────────────────
    def _resolve_title(self) -> str:
        return self._pending_title or self._last_heading or self.current_section or "Details"

    def heading(self, level: int, text: str):
        """A heading block. level<=1 -> section divider; level>=2 -> DEFER a content
        slide titled with this heading until its first body block arrives, so a heading
        followed by a table/flow/image never leaves an orphaned empty slide."""
        text = _txt(text, _MAX_TITLE_CHARS)
        if not text:
            return
        # Track the most recent heading of ANY level so tables/flows/images with no caption
        # are titled by their actual (sub)section, not a stale earlier level-1 section.
        self._last_heading = text
        if level <= 1:
            self.section_slide(text)         # resets body (sets current_section)
        else:
            self._reset_body()
            self._pending_title = text

    def _open_content_slide(self, title: Optional[str] = None):
        """Start a fresh content slide (title bar + footer) and reset the body cursor."""
        t = _txt(title if title is not None else self._resolve_title(), _MAX_TITLE_CHARS) or "Details"
        slide = self._new_slide()
        self._title_bar(slide, t)
        self._footer(slide)
        self._content_slide = slide
        self._content_y = _BODY_TOP
        self._current_title = t
        self._pending_title = ""             # consumed
        return slide

    # Back-compat alias: dedicated-slide blocks (table/flow/image) call this.
    def content_slide(self, title: str):
        return self._open_content_slide(title)

    def _ensure_room(self, needed: float):
        """Guarantee `needed` inches are free below the cursor on a content slide,
        spilling to a fresh continuation slide (same title) when the current one is full."""
        if self._content_slide is None:
            self._open_content_slide()
        elif self._content_y > _BODY_TOP + 0.01 and (self._content_y + needed) > _BODY_BOTTOM:
            self._open_content_slide(self._current_title)
        return self._content_slide

    # ── TITLE slide ─────────────────────────────────────────────────────────
    def title_slide(self, doc: Dict[str, Any]):
        slide = self._new_slide(bg=self.c_cover_bg)
        title = _txt(doc.get("title") or "Untitled Document", 200)
        subtitle = _txt(doc.get("subtitle") or "", 260)
        org = _txt(doc.get("org") or "", 120)
        doc_type = _txt(doc.get("doc_type") or "Document", 120)
        meta = "   ·   ".join([p for p in (org, doc_type, now_iso()[:10]) if p])

        hero = _decode_data_url(doc.get("cover_image_data_url"))
        placed_hero = False
        if hero is not None:
            try:
                placed_hero = self._place_cover_hero(slide, hero)
            except Exception as e:
                logger.warning("docgen pptx: cover hero failed: %s", e)
                placed_hero = False

        if placed_hero:
            # Text occupies the left ~55% over the cover background.
            text_w = _SLIDE_W * 0.55 - _MARGIN
            self._textbox(
                slide, title, _MARGIN, 2.25, text_w, 2.2,
                size=40, bold=True, color=self.c_cover_ink, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.BOTTOM, font=self.pal.get("font_bold"),
            )
            self._rect(slide, MSO_SHAPE.RECTANGLE, _MARGIN, 4.55, 2.2, 0.06, fill=self.c_cover_accent)
            if subtitle:
                self._textbox(
                    slide, subtitle, _MARGIN, 4.75, text_w, 1.1,
                    size=18, color=self.c_cover_ink, align=PP_ALIGN.LEFT,
                )
            if meta:
                self._textbox(
                    slide, meta, _MARGIN, 6.1, text_w, 0.6,
                    size=13, bold=True, color=self.c_cover_accent, align=PP_ALIGN.LEFT,
                )
        else:
            # Centered cover treatment (no hero / hero failed). Size the title to its
            # length and lay the underline / subtitle / meta out below its real height
            # so the accent rule never cuts through the title's descenders.
            tw = _SLIDE_W - 2 * _MARGIN
            tlen = len(title)
            t_size = 46 if tlen <= 42 else (40 if tlen <= 72 else 34)
            t_lines = _est_lines(title, tw - 0.4, t_size)
            t_h = max(1.0, t_lines * (t_size * 1.2 / 72.0))
            t_top = max(1.45, 3.5 - t_h / 2)     # keep the block visually centered
            self._textbox(
                slide, title, _MARGIN, t_top, tw, t_h + 0.1,
                size=t_size, bold=True, color=self.c_cover_ink, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.TOP, font=self.pal.get("font_bold"),
            )
            uy = t_top + t_h + 0.22
            self._rect(
                slide, MSO_SHAPE.RECTANGLE,
                _SLIDE_W / 2 - 1.3, uy, 2.6, 0.06, fill=self.c_cover_accent,
            )
            if subtitle:
                self._textbox(
                    slide, subtitle, _MARGIN, uy + 0.28, tw, 1.0,
                    size=20, color=self.c_cover_ink, align=PP_ALIGN.CENTER,
                )
            if meta:
                self._textbox(
                    slide, meta, _MARGIN, min(6.6, uy + 1.45), tw, 0.6,
                    size=13, bold=True, color=self.c_cover_accent, align=PP_ALIGN.CENTER,
                )
        self._reset_body()

    def _place_cover_hero(self, slide, data: bytes) -> bool:
        """Place the cover image as a right-half hero (or full-bleed if very wide). True if placed."""
        size = _image_size(data)
        if not size:
            return False
        w_px, h_px = size
        if w_px <= 0 or h_px <= 0:
            return False
        ratio = w_px / h_px

        if ratio >= 2.2:
            # Panoramic -> full-bleed band across the bottom 45% so the title sits above it.
            band_h = _SLIDE_H * 0.42
            band_top = _SLIDE_H - band_h
            self._add_picture_cover(slide, data, 0, band_top, _SLIDE_W, band_h, ratio)
            return True

        # Right-half hero panel.
        panel_left = _SLIDE_W * 0.55
        panel_w = _SLIDE_W - panel_left
        self._add_picture_cover(slide, data, panel_left, 0, panel_w, _SLIDE_H, ratio)
        return True

    def _add_picture_cover(self, slide, data, left, top, width, height, ratio):
        """Add a picture that COVERS the (left,top,width,height) box, center-cropped."""
        box_ratio = width / height if height else 1.0
        if ratio >= box_ratio:
            # image wider -> match height, overflow width (then clamp to box width)
            pic_h = height
            pic_w = height * ratio
        else:
            pic_w = width
            pic_h = width / ratio
        pic_left = left + (width - pic_w) / 2
        pic_top = top + (height - pic_h) / 2
        pic = slide.shapes.add_picture(
            BytesIO(data), Inches(pic_left), Inches(pic_top), Inches(pic_w), Inches(pic_h),
        )
        # Crop the overflow so it stays within the intended box.
        try:
            if pic_w > width:
                crop = (pic_w - width) / pic_w / 2
                pic.crop_left = crop
                pic.crop_right = crop
            if pic_h > height:
                crop = (pic_h - height) / pic_h / 2
                pic.crop_top = crop
                pic.crop_bottom = crop
            # Reset to the exact target box after cropping.
            pic.left = Inches(left)
            pic.top = Inches(top)
            pic.width = Inches(width)
            pic.height = Inches(height)
        except Exception:
            pass

    # ── AGENDA slide ─────────────────────────────────────────────────────────
    def agenda_slide(self, headings: List[Dict[str, Any]]):
        # Top-level sections only — a clean, numbered agenda that always fits.
        items = [h for h in headings if int(h.get("level", 1) or 1) <= 1][:_MAX_AGENDA_ITEMS]
        if not items:
            return
        slide = self._new_slide()
        self._title_bar(slide, "Agenda")
        # Scale font + spacing to the item count so the list never overflows the slide.
        n = len(items)
        if n <= 8:
            fsize, gap = 20, 10
        elif n <= 12:
            fsize, gap = 17, 7
        elif n <= 16:
            fsize, gap = 15, 5
        else:
            fsize, gap = 13, 3
        tf_box = slide.shapes.add_textbox(
            Inches(_BODY_LEFT), Inches(_BODY_TOP),
            Inches(_BODY_WIDTH), Inches(_BODY_HEIGHT),
        )
        tf = tf_box.text_frame
        tf.word_wrap = True
        num = 0
        first = True
        for h in items:
            text = _txt(h.get("text", ""), 110)
            if not text:
                continue
            num += 1
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(gap)
            r_num = p.add_run()
            r_num.text = f"{num:02d}   "
            r_num.font.size = Pt(fsize)
            r_num.font.bold = True
            r_num.font.color.rgb = self.c_accent
            r_main = p.add_run()
            r_main.text = text
            r_main.font.size = Pt(fsize)
            r_main.font.bold = True
            r_main.font.color.rgb = self.c_ink
        self._footer(slide)
        self._reset_body()

    # ── SECTION DIVIDER slide ────────────────────────────────────────────────
    def section_slide(self, text: str):
        self._section_index += 1
        slide = self._new_slide(bg=self.c_cover_bg)
        self.current_section = _txt(text, _MAX_TITLE_CHARS)
        # Big section number.
        self._textbox(
            slide, f"{self._section_index:02d}",
            _MARGIN, 1.7, _SLIDE_W - 2 * _MARGIN, 2.0,
            size=120, bold=True, color=self.c_cover_accent, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE, font=self.pal.get("font_bold"),
        )
        self._rect(slide, MSO_SHAPE.RECTANGLE, _MARGIN, 4.05, 2.4, 0.07, fill=self.c_cover_accent)
        self._textbox(
            slide, self.current_section,
            _MARGIN, 4.25, _SLIDE_W - 2 * _MARGIN, 1.8,
            size=36, bold=True, color=self.c_cover_ink, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, font=self.pal.get("font_bold"),
        )
        self._reset_body()

    # ── CONTENT body blocks (unified vertical flow) ───────────────────────────
    def add_paragraph(self, text: str):
        text = _txt(text, _MAX_PARA_CHARS)
        if not text:
            return
        eff_w = _BODY_WIDTH - 2 * _TB_INSET
        h = min(_est_text_height(text, eff_w, 16), _BODY_HEIGHT)
        self._ensure_room(h)
        self._textbox(
            self._content_slide, text,
            _BODY_LEFT, self._content_y, _BODY_WIDTH, h + 0.05,
            size=16, color=self.c_body, anchor=MSO_ANCHOR.TOP,
        )
        self._content_y += h + _PARA_GAP

    def _set_bullet(self, p, ordered: bool, start_at: int) -> None:
        """Apply a native PowerPoint hanging bullet (marL/indent) so wrapped lines align under
        the text instead of falling back to the left margin."""
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "274320")     # 0.30in body indent
        pPr.set("indent", "-228600")  # -0.25in hanging (marker sits left of the text)
        buClr = pPr.makeelement(qn("a:buClr"), {})
        buClr.append(pPr.makeelement(qn("a:srgbClr"), {"val": self._accent_hex}))
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
        if ordered:
            bu = pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod", "startAt": str(max(1, start_at))})
        else:
            bu = pPr.makeelement(qn("a:buChar"), {"char": "▪"})
        for el in (buClr, buFont, bu):
            pPr.append(el)

    def add_bullets(self, items: List[Any], ordered: bool = False):
        items = [_txt(x, _MAX_PARA_CHARS) for x in (items or [])]
        items = [x for x in items if x]
        if not items:
            return
        eff_w = _BODY_WIDTH - 2 * _TB_INSET - 0.35   # account for the bullet glyph indent
        i, n = 0, len(items)
        while i < n:
            first_h = _est_text_height(items[i], eff_w, 15) + 0.08
            self._ensure_room(min(first_h, _BODY_HEIGHT))
            avail = _BODY_BOTTOM - self._content_y
            group: List[Tuple[int, str]] = []
            gh = 0.0
            while i < n:
                ih = _est_text_height(items[i], eff_w, 15) + 0.08
                if group and (gh + ih) > avail:
                    break
                group.append((i, items[i]))
                gh += ih
                i += 1
            box = self._content_slide.shapes.add_textbox(
                Inches(_BODY_LEFT), Inches(self._content_y),
                Inches(_BODY_WIDTH), Inches(min(gh + 0.05, _BODY_HEIGHT)),
            )
            tf = box.text_frame
            tf.word_wrap = True
            for k, (orig_i, item) in enumerate(group):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.space_after = Pt(5)
                r = p.add_run()
                r.text = item
                r.font.size = Pt(15)
                r.font.color.rgb = self.c_body
                self._set_bullet(p, ordered, orig_i + 1)
            self._content_y += gh + _PARA_GAP

    # ── CALLOUT ───────────────────────────────────────────────────────────────
    def callout_shape(self, style: str, title: str, text: str):
        accent_hex, soft_bg_hex = _theme.CALLOUT_COLORS.get(style, _theme.CALLOUT_COLORS["info"])
        accent = _rgb(accent_hex)
        soft_bg = _rgb(soft_bg_hex)
        title = _txt(title or style.capitalize(), 120)
        text = _txt(text, _MAX_CALLOUT_CHARS)
        if not text and not title:
            return

        # Size the callout to its content so it never overlaps or overflows.
        eff_w = _BODY_WIDTH - 0.55           # internal L/R margins of the callout
        height = 0.30                        # top + bottom padding
        if title:
            height += _est_text_height(title, eff_w, 16)
        if text:
            height += _est_text_height(text, eff_w, 14) + (0.07 if title else 0)
        height = max(0.7, min(height, _BODY_HEIGHT))

        self._ensure_room(height)
        slide = self._content_slide
        top = self._content_y

        shape = self._rect(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            _BODY_LEFT, top, _BODY_WIDTH, height,
            fill=soft_bg, line=accent, line_w=1.5,
        )
        # Accent bar on the left edge for a stronger callout look.
        self._rect(slide, MSO_SHAPE.RECTANGLE, _BODY_LEFT, top, 0.09, height, fill=accent)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = title
        r0.font.bold = True
        r0.font.size = Pt(16)
        r0.font.color.rgb = accent
        if text:
            p1 = tf.add_paragraph()
            p1.space_before = Pt(4)
            r1 = p1.add_run()
            r1.text = text
            r1.font.size = Pt(14)
            r1.font.color.rgb = self.c_body
        self._content_y += height + _BLOCK_GAP

    # ── TABLE (own slide, native pptx table) ─────────────────────────────────
    def table_slide(self, columns: List[Any], rows: List[List[Any]], caption: str = ""):
        cols = [_txt(c, 60) for c in (columns or [])][:_MAX_TABLE_COLS]
        body_rows = [r for r in (rows or []) if isinstance(r, (list, tuple))]
        if not cols and body_rows:
            cols = [f"Column {i + 1}" for i in range(min(_MAX_TABLE_COLS, max(len(r) for r in body_rows)))]
        if not cols:
            return
        n_cols = len(cols)
        body_rows = body_rows[: _MAX_TABLE_ROWS - 1]
        norm_rows: List[List[str]] = []
        for r in body_rows:
            cells = [_txt(c, _MAX_CELL_CHARS) for c in list(r)[:n_cols]]
            cells += [""] * (n_cols - len(cells))
            norm_rows.append(cells)

        slide = self.content_slide(
            _txt(caption, _MAX_TITLE_CHARS) or self._pending_title or self._last_heading or self.current_section or "Table")
        n_rows = len(norm_rows) + 1   # +header
        top = _BODY_TOP
        height = min(_BODY_HEIGHT, 0.5 * n_rows + 0.2)
        gfx = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(_BODY_LEFT), Inches(top), Inches(_BODY_WIDTH), Inches(height),
        )
        table = gfx.table
        # Strip the built-in light style so our fills show.
        try:
            tbl = table._tbl
            tbl.firstRow = False
            tbl.horzBanding = False
        except Exception:
            pass

        def _style_cell(cell, text, *, header, zebra):
            cell.fill.solid()
            if header:
                cell.fill.fore_color.rgb = self.c_primary
            else:
                cell.fill.fore_color.rgb = self.c_panel_alt if zebra else self.c_panel
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            run.font.size = Pt(11 if header else 10)
            run.font.bold = bool(header)
            run.font.color.rgb = self.c_white if header else self.c_body

        for c, col in enumerate(cols):
            _style_cell(table.cell(0, c), col, header=True, zebra=False)
        for ri, row in enumerate(norm_rows, start=1):
            for c in range(n_cols):
                _style_cell(table.cell(ri, c), row[c] if c < len(row) else "",
                            header=False, zebra=(ri % 2 == 0))
        self._reset_body()

    # ── FLOW (own slide, chain of rounded rects + arrows) ─────────────────────
    def flow_slide(self, steps: List[Any], title: str = ""):
        steps = [_txt(models.flatten_step(s), _MAX_STEP_CHARS) for s in (steps or [])]
        steps = [s for s in steps if s][:_MAX_FLOW_STEPS]
        if not steps:
            return
        slide = self.content_slide(
            _txt(title, _MAX_TITLE_CHARS) or self._pending_title or self._last_heading or self.current_section or "Process Flow")
        n = len(steps)

        # Horizontal chain when few & short; otherwise a vertical chain.
        horizontal = n <= 4 and all(len(s) <= 38 for s in steps)
        if horizontal:
            self._flow_horizontal(slide, steps)
        else:
            self._flow_vertical(slide, steps)
        self._reset_body()

    def _flow_box(self, slide, step_no, text, left, top, width, height):
        shape = self._rect(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
            fill=self.c_panel, line=self.c_primary, line_w=1.25,
        )
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r_n = p.add_run()
        r_n.text = f"{step_no}  "
        r_n.font.size = Pt(14)
        r_n.font.bold = True
        r_n.font.color.rgb = self.c_accent
        r = p.add_run()
        r.text = text
        r.font.size = Pt(13)
        r.font.color.rgb = self.c_body

    def _flow_vertical(self, slide, steps):
        n = len(steps)
        area_top = _BODY_TOP + 0.1
        area_h = _BODY_HEIGHT - 0.2
        gap = 0.32
        box_h = min(0.95, (area_h - gap * (n - 1)) / n) if n else 0.95
        box_w = min(_BODY_WIDTH, 9.0)
        left = _BODY_LEFT + (_BODY_WIDTH - box_w) / 2
        y = area_top
        for i, step in enumerate(steps):
            self._flow_box(slide, i + 1, step, left, y, box_w, box_h)
            if i < n - 1:
                arrow = self._rect(
                    slide, MSO_SHAPE.DOWN_ARROW,
                    left + box_w / 2 - 0.12, y + box_h + 0.02,
                    0.24, max(0.12, gap - 0.06), fill=self.c_accent,
                )
            y += box_h + gap

    def _flow_horizontal(self, slide, steps):
        n = len(steps)
        gap = 0.4
        total_gap = gap * (n - 1)
        box_w = (_BODY_WIDTH - total_gap) / n
        box_h = 1.5
        y = _BODY_TOP + (_BODY_HEIGHT - box_h) / 2
        x = _BODY_LEFT
        for i, step in enumerate(steps):
            self._flow_box(slide, i + 1, step, x, y, box_w, box_h)
            if i < n - 1:
                arrow = self._rect(
                    slide, MSO_SHAPE.RIGHT_ARROW,
                    x + box_w + 0.04, y + box_h / 2 - 0.12,
                    max(0.12, gap - 0.08), 0.24, fill=self.c_accent,
                )
            x += box_w + gap

    # ── IMAGE (own slide) ─────────────────────────────────────────────────────
    def image_slide(self, block: Dict[str, Any]):
        caption = _txt(block.get("caption") or block.get("alt") or "", _MAX_CAPTION_CHARS)
        # Prefer a short slide title (section name); only use the caption as the title if it is
        # short. A long descriptive caption belongs UNDER the image, not as a wrapping title.
        short_caption = caption if 0 < len(caption) <= 48 else ""
        title = self._pending_title or self._last_heading or self.current_section or short_caption or "Figure"
        slide = self.content_slide(_txt(title, _MAX_TITLE_CHARS))
        data = _decode_data_url(block.get("data_url"))

        # Frame the image area in the body region.
        frame_left = _BODY_LEFT
        frame_top = _BODY_TOP + 0.1
        frame_w = _BODY_WIDTH
        frame_h = _BODY_HEIGHT - (0.5 if caption else 0.15)

        placed = False
        if data is not None:
            size = _image_size(data)
            if size:
                w_px, h_px = size
                if w_px > 0 and h_px > 0:
                    ratio = w_px / h_px
                    fit_w = frame_w
                    fit_h = fit_w / ratio
                    if fit_h > frame_h:
                        fit_h = frame_h
                        fit_w = fit_h * ratio
                    pic_left = frame_left + (frame_w - fit_w) / 2
                    pic_top = frame_top + (frame_h - fit_h) / 2
                    try:
                        slide.shapes.add_picture(
                            BytesIO(data),
                            Inches(pic_left), Inches(pic_top),
                            Inches(fit_w), Inches(fit_h),
                        )
                        placed = True
                    except Exception as e:
                        logger.warning("docgen pptx: image add_picture failed: %s", e)
                        placed = False

        if not placed:
            self._image_placeholder(slide, frame_left, frame_top, frame_w, frame_h,
                                    block.get("alt") or block.get("prompt") or caption)

        if caption:
            self._textbox(
                slide, caption,
                _BODY_LEFT, _BODY_BOTTOM - 0.35, _BODY_WIDTH, 0.35,
                size=11, color=self.c_muted, align=PP_ALIGN.CENTER,
            )
        self._reset_body()

    def _image_placeholder(self, slide, left, top, width, height, label):
        self._rect(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
            fill=self.c_panel, line=self.c_muted, line_w=1.0,
        )
        # Centered glyph.
        self._textbox(
            slide, "▣",
            left, top + height / 2 - 0.85, width, 0.9,
            size=44, color=self.c_muted, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        label = _txt(label or "Image", 180)
        if label:
            self._textbox(
                slide, label,
                left + 0.5, top + height / 2 + 0.1, width - 1.0, 0.9,
                size=13, color=self.c_muted, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.TOP,
            )

    # ── CLOSING slide ─────────────────────────────────────────────────────────
    def closing_slide(self, doc: Dict[str, Any]):
        slide = self._new_slide(bg=self.c_cover_bg)
        self._textbox(
            slide, "Thank you",
            _MARGIN, 2.6, _SLIDE_W - 2 * _MARGIN, 1.6,
            size=52, bold=True, color=self.c_cover_ink, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, font=self.pal.get("font_bold"),
        )
        self._rect(
            slide, MSO_SHAPE.RECTANGLE,
            _SLIDE_W / 2 - 1.3, 4.3, 2.6, 0.06, fill=self.c_cover_accent,
        )
        org = _txt(doc.get("org") or "", 120)
        if org:
            self._textbox(
                slide, org,
                _MARGIN, 4.55, _SLIDE_W - 2 * _MARGIN, 0.7,
                size=18, color=self.c_cover_ink, align=PP_ALIGN.CENTER,
            )

    # ── single-slide FLYER ────────────────────────────────────────────────────
    def flyer_slide(self, doc: Dict[str, Any]):
        """A one-slide promotional flyer: top hero band, headline, tagline, punchy benefit
        list, and a bottom call-to-action band."""
        slide = self._new_slide(bg=self.c_cover_bg)
        W, H = _SLIDE_W, _SLIDE_H
        band_h = H * 0.42
        margin = 0.85
        cw = W - 2 * margin
        cta_h = 1.1

        # Hero band (cover-fit) across the top, else a solid primary_dark band.
        data = _decode_data_url(doc.get("cover_image_data_url"))
        placed = False
        if data is not None:
            size = _image_size(data)
            if size and size[0] > 0 and size[1] > 0:
                try:
                    self._add_picture_cover(slide, data, 0, 0, W, band_h, size[0] / size[1])
                    placed = True
                except Exception:
                    placed = False
        if not placed:
            self._rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, band_h, fill=self.c_primary_dark)
        self._rect(slide, MSO_SHAPE.RECTANGLE, 0, band_h - 0.03, W, 0.06, fill=self.c_cover_accent)

        # Headline below the hero, with an accent tab.
        self._rect(slide, MSO_SHAPE.RECTANGLE, margin, band_h + 0.18, 1.4, 0.06, fill=self.c_accent)
        self._textbox(
            slide, _txt(doc.get("title") or "Untitled", 120),
            margin, band_h + 0.22, cw, 0.9,
            size=31, bold=True, color=self.c_cover_ink, anchor=MSO_ANCHOR.TOP,
            font=self.pal.get("font_bold"),
        )

        tagline, bullets, cta_title, cta_text = _flyer_extract(doc)
        y = band_h + 1.02
        if tagline:
            self._textbox(slide, _txt(tagline, 240), margin, y, cw, 0.6,
                          size=15, color=self.c_cover_ink, anchor=MSO_ANCHOR.TOP)
            y += 0.56
        for it in bullets:
            if y > H - cta_h - 0.28:
                break
            self._rect(slide, MSO_SHAPE.RECTANGLE, margin, y + 0.07, 0.13, 0.13, fill=self.c_accent)
            self._textbox(slide, _txt(it, 180), margin + 0.34, y - 0.04, cw - 0.34, 0.5,
                          size=13.5, color=self.c_cover_ink, anchor=MSO_ANCHOR.TOP)
            y += 0.46

        # Call-to-action band.
        self._rect(slide, MSO_SHAPE.RECTANGLE, 0, H - cta_h, W, cta_h, fill=self.c_primary)
        self._rect(slide, MSO_SHAPE.RECTANGLE, 0, H - cta_h, W, 0.05, fill=self.c_cover_accent)
        if cta_title:
            self._textbox(slide, _txt(cta_title, 120), margin, H - cta_h + 0.24, cw, 0.5,
                          size=18, bold=True, color=self.c_white, align=PP_ALIGN.CENTER,
                          font=self.pal.get("font_bold"))
        if cta_text:
            self._textbox(slide, _txt(cta_text, 220), margin, H - cta_h + 0.78, cw, 0.45,
                          size=12.5, color=self.c_white, align=PP_ALIGN.CENTER)


def _collect_headings(blocks: List[Any]) -> List[Dict[str, Any]]:
    """Level-1 (and nested level-2) headings for the agenda."""
    out: List[Dict[str, Any]] = []
    for b in blocks:
        if not isinstance(b, dict) or str(b.get("type", "")).lower() != "heading":
            continue
        try:
            level = int(b.get("level", 2) or 2)
        except (TypeError, ValueError):
            level = 2
        text = _txt(b.get("text", ""), 120)
        if text and level <= 2:
            out.append({"level": level, "text": text})
    # Only show an agenda if there's at least one top-level section.
    if not any(h["level"] <= 1 for h in out):
        return []
    return out


def render_pptx(doc: Dict[str, Any]) -> bytes:
    """Render a (normalized) document dict to PPTX bytes. Never raises on bad input."""
    if not isinstance(doc, dict):
        doc = {}

    try:
        pal = _theme.get_theme(doc.get("theme"))
    except Exception:
        pal = _theme.get_theme(None)

    deck = _Deck(pal)
    deck.org = _txt(doc.get("org") or "", 80)

    # Single-slide promotional flyer layout.
    if str(doc.get("layout", "")).lower() == "flyer":
        try:
            deck.flyer_slide(doc)
            buf = BytesIO()
            deck.prs.save(buf)
            return buf.getvalue()
        except Exception as e:
            logger.warning("docgen pptx: flyer render failed (%s); falling back to document", e)
            deck = _Deck(pal)
            deck.org = _txt(doc.get("org") or "", 80)

    # Always emit the title slide.
    try:
        deck.title_slide(doc)
    except Exception as e:
        logger.warning("docgen pptx: title slide failed: %s", e)

    blocks = doc.get("blocks")
    if not isinstance(blocks, list):
        blocks = []

    # Agenda (best-effort) from top-level headings.
    try:
        headings = _collect_headings(blocks)
        if headings:
            deck.agenda_slide(headings)
    except Exception as e:
        logger.warning("docgen pptx: agenda failed: %s", e)

    for block in blocks:
        if deck.full():
            break
        if not isinstance(block, dict):
            continue
        btype = str(block.get("type", "")).strip().lower()
        try:
            if btype == "heading":
                try:
                    level = int(block.get("level", 2))
                except (TypeError, ValueError):
                    level = 2
                deck.heading(level, block.get("text") or "")

            elif btype == "paragraph":
                deck.add_paragraph(block.get("text") or "")

            elif btype == "bullets":
                deck.add_bullets(block.get("items") or [], bool(block.get("ordered")))

            elif btype == "table":
                deck.table_slide(
                    block.get("columns") or [],
                    block.get("rows") or [],
                    _txt(block.get("caption") or "", _MAX_TITLE_CHARS),
                )

            elif btype == "flow":
                deck.flow_slide(block.get("steps") or [], _txt(block.get("title") or "", _MAX_TITLE_CHARS))

            elif btype == "callout":
                deck.callout_shape(
                    str(block.get("style", "info")).strip().lower(),
                    block.get("title") or "",
                    block.get("text") or "",
                )

            elif btype == "image":
                deck.image_slide(block)

            elif btype in ("divider", "pagebreak"):
                # Force subsequent body text onto a fresh slide.
                deck._reset_body()
            # unknown/malformed block types are silently skipped
        except Exception as e:
            logger.warning("docgen pptx: block %r failed: %s", btype, e)
            deck._reset_body()
            continue

    # Always emit a closing slide (within the cap).
    if not deck.full():
        try:
            deck.closing_slide(doc)
        except Exception as e:
            logger.warning("docgen pptx: closing slide failed: %s", e)

    buf = BytesIO()
    deck.prs.save(buf)
    return buf.getvalue()
