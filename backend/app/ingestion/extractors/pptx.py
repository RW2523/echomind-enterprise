"""
PPTX extractor: extract text per slide for chunking.
"""
from __future__ import annotations
from typing import List, Dict, Any
from io import BytesIO
from pptx import Presentation


def extract_pptx(data: bytes) -> List[Dict[str, Any]]:
    """
    Extract text per slide. Returns list of {"text", "page_start", "page_end", "section_path"}.
    """
    blocks = []
    prs = Presentation(BytesIO(data))
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        text = "\n".join(parts).strip()
        if text:
            blocks.append({
                "text": text,
                "page_start": i + 1,
                "page_end": i + 1,
                "section_path": f"slide_{i + 1}",
            })
    return blocks
